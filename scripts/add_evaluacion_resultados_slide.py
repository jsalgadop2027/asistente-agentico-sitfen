"""Completa exclusivamente la diapositiva 34 de la presentación final de SITFEN.

La diapositiva ya existe y está vacía entre los separadores de los capítulos X y XI.
El resto del contenido no se inserta, elimina, reordena ni edita.
"""
from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
INPUT = ROOT / "entregables" / "SITFEN_IA_Exposición_Jurado_Final.pptx"
OUTPUT = ROOT / "entregables" / "SITFEN_IA_Exposición_Jurado_Final_v2.pptx"
SLIDE_INDEX = 33  # diapositiva 34, índice base cero

FONT = "Space Grotesk"
FONT_MEDIUM = "Space Grotesk Medium"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value.lstrip("#").upper())


NAVY = rgb("082345")
DARK = rgb("4D4B4C")
TEXT = rgb("26313D")
MUTED = rgb("6B6B6B")
WHITE = rgb("FFFFFF")
LIGHT = rgb("E4E7EA")
CYAN = rgb("00C8DE")
BLUE = rgb("1168C4")
TEAL = rgb("07969D")
PURPLE = rgb("6A43B5")
AMBER = rgb("DE9400")
RED = rgb("D23B43")
GREEN = rgb("2F7D45")


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 12,
    color: RGBColor = TEXT,
    bold: bool = False,
    font: str = FONT,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.01,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    p.line_spacing = 1.0
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor = WHITE,
    line: RGBColor | None = LIGHT,
    radius: bool = True,
    line_width: float = 1.0,
):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width)
    return shape


def add_card(slide, x: float, y: float, w: float, h: float, *, fill: RGBColor, line: RGBColor):
    add_rect(slide, x + 0.035, y + 0.045, w, h, fill=rgb("C9CDD1"), line=None)
    return add_rect(slide, x, y, w, h, fill=fill, line=line, line_width=1.3)


def add_metric_card(
    slide,
    x: float,
    *,
    label: str,
    value: str,
    note: str,
    accent: RGBColor,
    dark: bool = False,
    value_size: float = 29,
):
    y, w, h = 1.50, 2.55, 1.30
    fill = NAVY if dark else WHITE
    add_card(slide, x, y, w, h, fill=fill, line=accent)
    title_color = CYAN if dark else accent
    body_color = WHITE if dark else TEXT
    note_color = WHITE if dark else MUTED
    add_text(slide, label.upper(), x + 0.17, y + 0.15, w - 0.34, 0.19, size=7.7, color=title_color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, value, x + 0.17, y + 0.37, w - 0.34, 0.44, size=value_size, color=body_color if dark else accent, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, note, x + 0.18, y + 0.87, w - 0.36, 0.30, size=7.3, color=note_color, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def add_metric_bar(slide, label: str, value: float, y: float, color: RGBColor) -> None:
    add_text(slide, label, 0.91, y - 0.01, 1.21, 0.17, size=7.4, color=TEXT, bold=True)
    add_rect(slide, 2.13, y, 2.02, 0.13, fill=LIGHT, line=None, radius=True)
    add_rect(slide, 2.13, y, 2.02 * value, 0.13, fill=color, line=None, radius=True)
    add_text(slide, f"{value:.4f}", 4.21, y - 0.04, 0.52, 0.20, size=7.6, color=color, bold=True, align=PP_ALIGN.RIGHT)


def add_check(slide, y: float, title: str, body: str, color: RGBColor) -> None:
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.66), Inches(y), Inches(0.28), Inches(0.28))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_text(slide, "✓", 5.66, y, 0.28, 0.27, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)
    add_text(slide, title, 6.06, y - 0.01, 1.15, 0.18, size=7.6, color=color, bold=True)
    add_text(slide, body, 7.17, y - 0.01, 2.10, 0.28, size=7.3, color=TEXT)


def build_slide(slide) -> None:
    if len(slide.shapes):
        raise RuntimeError("La diapositiva 34 ya contiene elementos; se cancela para no sobrescribir contenido.")

    add_text(
        slide,
        "Alta fidelidad; recuperación por fortalecer",
        0.67,
        0.88,
        8.45,
        0.43,
        size=25.0,
        color=DARK,
        bold=True,
        font=FONT_MEDIUM,
        valign=MSO_ANCHOR.MIDDLE,
    )
    underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.67), Inches(1.27), Inches(0.74), Inches(0.035))
    underline.fill.solid()
    underline.fill.fore_color.rgb = CYAN
    underline.line.fill.background()

    add_metric_card(
        slide,
        0.67,
        label="Fidelidad",
        value="0.94",
        note="RAGAS 0.9375 · DeepEval 0.9458\ntríada propia 1.0000",
        accent=CYAN,
        dark=True,
    )
    add_metric_card(
        slide,
        3.41,
        label="Relevancia",
        value="0.91",
        note="DeepEval 0.9062 · meta ≥0.80\nRAGAS penaliza la abstención honesta",
        accent=PURPLE,
    )
    add_metric_card(
        slide,
        6.15,
        label="Costo variable",
        value="$0.0137",
        note="US$/consulta · techo $0.015\n8.7% por debajo del KPI-2",
        accent=TEAL,
        value_size=25,
    )

    add_card(slide, 0.67, 3.04, 4.20, 1.42, fill=WHITE, line=AMBER)
    add_text(slide, "RECUPERACIÓN · DIAGNÓSTICO ACCIONABLE", 0.89, 3.20, 3.75, 0.18, size=7.7, color=AMBER, bold=True)
    add_metric_bar(slide, "Context precision", 0.6424, 3.55, BLUE)
    add_metric_bar(slide, "Context relevancy", 0.6780, 3.84, TEAL)
    add_metric_bar(slide, "Context recall", 0.3333, 4.13, RED)
    add_text(slide, "Recall bajo = vacíos del corpus + desalineamiento del golden.", 0.92, 4.31, 3.68, 0.14, size=6.2, color=MUTED, align=PP_ALIGN.CENTER)

    add_card(slide, 5.10, 3.04, 4.15, 1.42, fill=WHITE, line=TEAL)
    add_text(slide, "OPERACIÓN E2E VERIFICADA", 5.35, 3.20, 3.60, 0.18, size=7.7, color=TEAL, bold=True)
    add_check(slide, 3.52, "CANAL", "WhatsApp texto/voz · memoria · fuentes", BLUE)
    add_check(slide, 3.87, "CONTROL", "Admin UI · EICAR · macros · duplicados", RED)
    add_check(slide, 4.22, "NUBE", "GCP serverless · health checks · CI en verde", GREEN)

    add_rect(slide, 0.67, 4.68, 8.58, 0.37, fill=CYAN, line=None, radius=True)
    add_text(slide, "La respuesta está validada; ahora toca ampliar evidencia sin sacrificar abstención.", 0.90, 4.75, 8.12, 0.19, size=9.2, color=DARK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    add_text(slide, "Fuente: Informe SITFEN v12 · Cap. X · Tabla 14 · corrida formal n=45", 0.67, 5.23, 5.85, 0.14, size=6.3, color=rgb("808080"), margin=0)

    try:
        slide.notes_slide.notes_text_frame.text = (
            "Lectura sugerida: la fidelidad de 0.94 confirma la hipótesis técnica de RAG con abstención. "
            "La relevancia DeepEval supera la meta; RAGAS reduce su puntuación cuando el asistente se abstiene "
            "ante vacíos del corpus. El punto débil medible es la recuperación, especialmente el context recall, "
            "por vacíos documentales y desalineamiento del golden dataset. Operativamente, el flujo E2E quedó "
            "verificado en GCP y el costo de US$0.0137 se mantiene 8.7% bajo el tope del KPI-2."
        )
    except Exception:
        pass


def main() -> None:
    if not INPUT.exists():
        raise FileNotFoundError(INPUT)
    shutil.copy2(INPUT, OUTPUT)
    prs = Presentation(OUTPUT)
    if len(prs.slides) != 48:
        raise RuntimeError(f"Se esperaban 48 diapositivas y se encontraron {len(prs.slides)}.")
    build_slide(prs.slides[SLIDE_INDEX])
    prs.core_properties.title = "SITFEN — Presentación final ante jurado"
    prs.core_properties.subject = "Presentación final con Evaluación de Resultados del Capítulo X"
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    main()
