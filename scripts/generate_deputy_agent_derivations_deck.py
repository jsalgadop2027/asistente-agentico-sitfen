from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "entregables" / "Presentacion_Diputado_Derivaciones_Agente_IA.pptx"

NAVY = RGBColor(9, 38, 76)
BLUE = RGBColor(19, 104, 170)
CYAN = RGBColor(0, 157, 184)
TEAL = RGBColor(0, 145, 135)
GREEN = RGBColor(37, 146, 104)
AMBER = RGBColor(225, 145, 31)
RED = RGBColor(190, 58, 66)
PURPLE = RGBColor(100, 77, 157)
INK = RGBColor(44, 55, 68)
MUTED = RGBColor(105, 116, 128)
LIGHT = RGBColor(222, 231, 239)
PALE = RGBColor(239, 246, 250)
WHITE = RGBColor(255, 255, 255)
BG = RGBColor(248, 251, 253)
FONT = "Aptos"


def add_text(slide, value, x, y, w, h, size=14, color=INK, bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = p.space_after = Pt(0)
    run = p.add_run()
    run.text = value
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_shape(slide, kind, x, y, w, h, fill=WHITE, line=LIGHT, radius=True):
    if kind == "rect":
        kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def add_line(slide, x1, y1, x2, y2, color=BLUE, width=2, dash=None):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    if dash:
        line.line.dash_style = dash
    return line


def add_arrow(slide, x, y, w, h, color=BLUE, direction="right"):
    kinds = {"right": MSO_SHAPE.RIGHT_ARROW, "left": MSO_SHAPE.LEFT_ARROW,
             "down": MSO_SHAPE.DOWN_ARROW, "up": MSO_SHAPE.UP_ARROW}
    return add_shape(slide, kinds[direction], x, y, w, h, color, None, False)


def icon(slide, label, x, y, color=BLUE, d=0.55):
    add_shape(slide, MSO_SHAPE.OVAL, x, y, d, d, color, None)
    add_text(slide, label, x, y + 0.05, d, d - 0.08, 12, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)


def card(slide, x, y, w, h, title, body, accent=BLUE, badge=None):
    add_shape(slide, "rect", x, y, w, h, WHITE, LIGHT)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.08, h, accent, None, False)
    if badge:
        icon(slide, badge, x + 0.20, y + 0.20, accent, 0.44)
        tx = x + 0.78
    else:
        tx = x + 0.22
    add_text(slide, title, tx, y + 0.18, w - (tx - x) - 0.18, 0.30, 13, NAVY, True)
    add_text(slide, body, x + 0.22, y + 0.62, w - 0.40, h - 0.74, 9.6, MUTED)


def base(prs, title, section, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.15, NAVY, None, False)
    add_text(slide, section.upper(), 0.58, 0.34, 1.40, 0.24, 10, CYAN, True)
    add_text(slide, title, 1.55, 0.27, 11.10, 0.42, 22, NAVY, True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.58, 0.82, 0.90, 0.045, CYAN, None, False)
    if subtitle:
        add_text(slide, subtitle, 0.58, 0.95, 12.1, 0.30, 10.5, MUTED)
    add_text(slide, "SITFEN | Asistente agéntico para derivaciones y escalamiento", 0.58, 7.16, 7.3, 0.18, 7.5, MUTED)
    add_text(slide, str(len(prs.slides)).zfill(2), 12.25, 7.13, 0.50, 0.20, 8, MUTED, True, PP_ALIGN.RIGHT)
    return slide


def pill(slide, value, x, y, w, fill=PALE, color=NAVY):
    add_shape(slide, "rect", x, y, w, 0.30, fill, None)
    add_text(slide, value, x + 0.05, y + 0.03, w - 0.10, 0.22, 8.5, color, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Executive framing
    s = base(prs, "Del mensaje ciudadano a la acción pública coordinada", "01 | Propósito",
             "Un asistente que comprende la necesidad, identifica el destino correcto y deja una traza auditable.")
    add_text(s, "El valor no está solo en responder.\nEstá en activar al Estado correcto, a tiempo y con contexto.",
             0.75, 1.55, 5.2, 1.05, 23, NAVY, True)
    add_text(s, "El ciudadano conserva un canal simple: WhatsApp. El sistema convierte una conversación en una solicitud estructurada, priorizada y derivable.",
             0.78, 2.78, 4.85, 0.70, 12, MUTED)
    # Main flow
    stages = [(0.75, "Ciudadano", "WhatsApp\ntexto, voz o imagen", GREEN, "1"),
              (3.10, "Asistente", "Entiende, busca\ny decide", PURPLE, "2"),
              (5.45, "Derivación", "Entidad + urgencia\n+ confirmación", AMBER, "3"),
              (7.80, "Entidad pública", "Recibe contexto\ny prioridad", BLUE, "4"),
              (10.15, "Seguimiento", "Handoff, estado\ny feedback", TEAL, "5")]
    for i, (x, title, body, color, num) in enumerate(stages):
        add_shape(s, "rect", x, 4.35, 1.78, 1.42, WHITE, LIGHT)
        icon(s, num, x + 0.16, 4.56, color, 0.42)
        add_text(s, title, x + 0.68, 4.58, 0.96, 0.25, 11.5, NAVY, True)
        add_text(s, body, x + 0.18, 5.12, 1.42, 0.48, 9, MUTED, False, PP_ALIGN.CENTER)
        if i < len(stages) - 1:
            add_arrow(s, x + 1.88, 4.92, 0.35, 0.22, color)
    pill(s, "Menos fricción", 0.78, 6.28, 1.50, PALE, BLUE)
    pill(s, "Más focalización", 2.42, 6.28, 1.62, PALE, TEAL)
    pill(s, "Decisiones trazables", 4.18, 6.28, 1.80, PALE, PURPLE)

    # 2. Architecture
    s = base(prs, "El asistente es un sistema de componentes, no un chatbot aislado", "02 | Arquitectura",
             "La orquestación agéntica conecta canales, conocimiento, decisiones y acciones con controles explícitos.")
    # central agent
    add_shape(s, MSO_SHAPE.OVAL, 5.30, 2.45, 2.55, 2.05, NAVY, CYAN)
    add_text(s, "AGENTE\nREACT", 5.67, 2.78, 1.80, 0.56, 22, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    add_text(s, "Orquesta herramientas\ny decide el siguiente paso", 5.60, 3.58, 1.95, 0.42, 9.5, RGBColor(190, 226, 236), False, PP_ALIGN.CENTER)
    blocks = [(0.75, 1.62, "Canales", "WhatsApp · Web\ntexto · voz · imagen", GREEN, "A"),
              (0.75, 4.55, "Memoria + RAG", "Fuentes oficiales\ncontexto por usuario", BLUE, "B"),
              (9.95, 1.62, "Clasificación", "Intención · sentir\nurgencia · entidad", AMBER, "C"),
              (9.95, 4.55, "Acciones", "Derivación · handoff\nreportes · alertas", TEAL, "D")]
    for x, y, title, body, color, mark in blocks:
        card(s, x, y, 2.62, 1.32, title, body, color, mark)
    add_line(s, 3.38, 2.28, 5.30, 3.00, GREEN, 2)
    add_line(s, 3.38, 5.20, 5.30, 3.95, BLUE, 2)
    add_line(s, 7.85, 3.00, 9.95, 2.28, AMBER, 2)
    add_line(s, 7.85, 3.95, 9.95, 5.20, TEAL, 2)
    add_text(s, "Ventaja institucional", 4.50, 5.18, 4.35, 0.28, 12, NAVY, True, PP_ALIGN.CENTER)
    add_text(s, "Cada componente puede observarse, evaluarse y mejorar sin perder el control del flujo.", 4.15, 5.52, 5.00, 0.42, 10.5, MUTED, False, PP_ALIGN.CENTER)
    pill(s, "Modular", 4.75, 6.30, 1.10, PALE, BLUE)
    pill(s, "Auditable", 6.02, 6.30, 1.20, PALE, PURPLE)
    pill(s, "Escalable", 7.38, 6.30, 1.15, PALE, TEAL)

    # 3. Derivation flow
    s = base(prs, "Derivación inteligente: una decisión antes de enviar", "03 | Flujo operativo",
             "La acción irreversible no queda a criterio del modelo: exige reglas, confirmación y un destino controlado.")
    steps = [(0.72, "1", "Detectar", "Pedido, reclamo,\nriesgo u objetivo", PURPLE),
             (2.82, "2", "Comprender", "Resumen del caso\ncon fuentes y contexto", BLUE),
             (4.92, "3", "Priorizar", "Urgencia: baja, media,\nalta o crítica", RED),
             (7.02, "4", "Clasificar", "Catálogo cerrado\nde 10 entidades", AMBER),
             (9.12, "5", "Confirmar", "El ciudadano valida\nel texto real", GREEN),
             (11.22, "6", "Derivar", "Fan-out deduplicado\ny asunto prioritario", TEAL)]
    for i, (x, num, title, body, color) in enumerate(steps):
        add_shape(s, "rect", x, 2.02, 1.48, 1.58, WHITE, LIGHT)
        icon(s, num, x + 0.48, 2.20, color, 0.46)
        add_text(s, title, x + 0.12, 2.82, 1.24, 0.24, 10.8, NAVY, True, PP_ALIGN.CENTER)
        add_text(s, body, x + 0.12, 3.15, 1.24, 0.34, 8.4, MUTED, False, PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            add_arrow(s, x + 1.57, 2.68, 0.40, 0.20, color)
    add_shape(s, "rect", 0.72, 4.28, 5.72, 1.42, PALE, LIGHT)
    add_text(s, "¿Qué recibe la entidad?", 0.98, 4.52, 2.42, 0.28, 13, NAVY, True)
    add_text(s, "Resumen del caso · entidad sugerida · nivel de urgencia · evidencia consultada · datos de contacto · consentimiento", 0.98, 4.92, 5.00, 0.48, 10.4, INK)
    add_shape(s, "rect", 6.82, 4.28, 5.78, 1.42, WHITE, LIGHT)
    add_text(s, "Destinos posibles", 7.08, 4.52, 2.20, 0.28, 13, NAVY, True)
    entities = ["SENASA", "SENAMHI / ENFEN", "INDECI", "MIDAGRI", "CITE / RedCITE", "PROMPERÚ", "SUNAT", "Gobierno Regional", "Municipalidad"]
    for i, entity in enumerate(entities):
        pill(s, entity, 7.08 + (i % 3) * 1.70, 4.94 + (i // 3) * 0.34, 1.55, PALE, BLUE if i % 2 else TEAL)
    add_text(s, "Resultado: la solicitud llega al responsable correcto, con menos rebotes y mejor priorización.", 1.15, 6.24, 11.10, 0.36, 13, NAVY, True, PP_ALIGN.CENTER)

    # 4. Human oversight and safety
    s = base(prs, "Escalamiento humano: la IA sabe cuándo ceder el control", "04 | Gobernanza",
             "El sistema combina autonomía operativa con supervisión humana para casos sensibles, ambiguos o urgentes.")
    add_shape(s, MSO_SHAPE.OVAL, 5.36, 2.05, 2.30, 2.30, NAVY, CYAN)
    add_text(s, "HUMANO\nEN EL\nCIRCUITO", 5.70, 2.58, 1.62, 0.80, 19, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    add_text(s, "Decide · valida · retoma", 5.55, 3.72, 1.92, 0.25, 9.5, RGBColor(190, 226, 236), False, PP_ALIGN.CENTER)
    cards = [(0.80, 1.82, "Handoff en vivo", "Registra el caso en Firestore y avisa al equipo para que una persona retome la conversación.", TEAL, "1"),
             (0.80, 4.48, "Urgencia visible", "Los casos altos o críticos se etiquetan para triage prioritario.", RED, "2"),
             (9.93, 1.82, "Confirmación explícita", "Derivar requiere una afirmación del usuario en el turno actual.", GREEN, "3"),
             (9.93, 4.48, "Trazabilidad", "Se conserva el contexto, la entidad, la decisión y el resultado para auditoría.", PURPLE, "4")]
    for x, y, title, body, color, mark in cards:
        card(s, x, y, 2.62, 1.28, title, body, color, mark)
    add_line(s, 3.42, 2.50, 5.36, 2.80, TEAL, 2)
    add_line(s, 3.42, 5.08, 5.36, 3.62, RED, 2)
    add_line(s, 7.66, 2.80, 9.93, 2.50, GREEN, 2)
    add_line(s, 7.66, 3.62, 9.93, 5.08, PURPLE, 2)
    add_text(s, "Esto protege al ciudadano y a la institución: velocidad para lo rutinario, criterio humano para lo que importa.", 1.30, 6.30, 10.70, 0.34, 13, NAVY, True, PP_ALIGN.CENTER)

    # 5. Value and ask
    s = base(prs, "Qué gana el Estado y cómo se mide", "05 | Impacto",
             "La propuesta convierte conversaciones dispersas en coordinación interinstitucional medible.")
    metrics = [(0.78, "01", "Acceso", "Un canal cotidiano para iniciar una solicitud sin conocer la estructura del Estado.", GREEN),
               (3.30, "02", "Focalización", "La derivación se dirige por dominio, urgencia y contexto; no por rebote manual.", BLUE),
               (5.82, "03", "Respuesta", "Las alertas críticas y los handoffs reducen el tiempo hasta una atención humana.", RED),
               (8.34, "04", "Evidencia", "Cada caso deja resumen, destino, prioridad, consentimiento y seguimiento.", PURPLE),
               (10.86, "05", "Aprendizaje", "Los cierres y reportes diarios muestran dónde se concentran los dolores ciudadanos.", TEAL)]
    for x, num, title, body, color in metrics:
        add_shape(s, "rect", x, 1.70, 2.18, 2.55, WHITE, LIGHT)
        icon(s, num, x + 0.80, 1.98, color, 0.56)
        add_text(s, title, x + 0.18, 2.75, 1.82, 0.30, 13.5, NAVY, True, PP_ALIGN.CENTER)
        add_text(s, body, x + 0.18, 3.22, 1.82, 0.70, 9.4, MUTED, False, PP_ALIGN.CENTER)
    add_shape(s, "rect", 0.78, 4.72, 7.26, 1.08, NAVY, None)
    add_text(s, "Indicadores para gestión pública", 1.08, 4.95, 2.50, 0.25, 12.5, WHITE, True)
    add_text(s, "% de casos con destino correcto  ·  tiempo a primera atención  ·  % de confirmaciones  ·  casos críticos atendidos  ·  tasa de cierre", 1.08, 5.32, 6.48, 0.28, 9.5, RGBColor(214, 233, 242))
    add_shape(s, "rect", 8.42, 4.72, 4.12, 1.08, PALE, LIGHT)
    add_text(s, "Decisión solicitada", 8.72, 4.95, 1.90, 0.25, 12.5, NAVY, True)
    add_text(s, "Impulsar un piloto interinstitucional con reglas de atención, correos verificados y responsables por entidad.", 8.72, 5.30, 3.48, 0.34, 9.5, INK)
    add_text(s, "Una puerta de entrada simple para el ciudadano. Una red de respuesta más inteligente para el Estado.", 1.10, 6.35, 11.15, 0.36, 16, NAVY, True, PP_ALIGN.CENTER)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build_deck()
