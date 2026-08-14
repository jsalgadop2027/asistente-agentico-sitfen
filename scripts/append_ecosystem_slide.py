from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "entregables" / "Presentacion_Diputado_Derivaciones_Agente_IA.pptx"
IMAGE = Path(r"E:\Mi perfil\Mis documentos\1. Actualización\UTEC\Maestría de CD & IA\Ciclo VI\Capstone Project II\Arquitectura del Agente\Ecosistema agéntico.drawio.png")


def append_slide():
    prs = Presentation(PPTX)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

    with Image.open(IMAGE) as image:
        image_ratio = image.width / image.height

    slide_ratio = prs.slide_width / prs.slide_height
    if image_ratio >= slide_ratio:
        width = prs.slide_width
        height = int(width / image_ratio)
        left = 0
        top = int((prs.slide_height - height) / 2)
    else:
        height = prs.slide_height
        width = int(height * image_ratio)
        top = 0
        left = int((prs.slide_width - width) / 2)

    slide.shapes.add_picture(str(IMAGE), left, top, width=width, height=height)
    prs.save(PPTX)
    print(f"{PPTX}\nslides={len(prs.slides)}")


if __name__ == "__main__":
    append_slide()
