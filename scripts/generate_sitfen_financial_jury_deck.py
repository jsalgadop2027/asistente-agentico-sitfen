"""Genera el módulo ejecutivo del Capítulo XI para la defensa de SITFEN.

Salida principal:
  - entregables/SITFEN_Capitulo_XI_Solo_Slides_Financieros.pptx

Las cifras se leen directamente de evaluacion_financiera/model_output.json.
"""
from __future__ import annotations

import json
import math
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
MODEL_PATH = ROOT / "evaluacion_financiera" / "model_output.json"
SOURCE_DECK = ROOT / "entregables" / "Exposición_SITFEN_Jurado_IA_v3.pptx"
OUT_MODULE = ROOT / "entregables" / "SITFEN_Capitulo_XI_Solo_Slides_Financieros.pptx"
OUT_INTEGRATED = ROOT / "entregables" / "Exposición_SITFEN_Jurado_IA_v5.pptx"
CHART_DIR = ROOT / "_work_presentation" / "financial_jury_charts"

FONT = "Space Grotesk"
FONT_MEDIUM = "Space Grotesk Medium"


def rgb(hex_value: str) -> RGBColor:
    return RGBColor.from_string(hex_value.replace("#", "").upper())


NAVY = rgb("082345")
INK = rgb("26313D")
WHITE = rgb("FFFFFF")
OFFWHITE = rgb("F6F8FA")
LIGHT = rgb("E2E8EE")
MID = rgb("AEB9C3")
MUTED = rgb("647383")
CYAN = rgb("00D8EE")
CYAN_DARK = rgb("0091A5")
BLUE = rgb("0B61B9")
TEAL = rgb("0A9196")
GREEN = rgb("2F7D45")
AMBER = rgb("D69200")
RED = rgb("C83942")
PURPLE = rgb("6945A5")

SLIDE_W = 10.0
SLIDE_H = 5.625


def money(value: float, decimals: int = 0, plus: bool = False) -> str:
    if value < 0:
        return f"-US$ {abs(value):,.{decimals}f}"
    sign = "+" if plus and value > 0 else ""
    return f"{sign}US$ {value:,.{decimals}f}"


def pct(value: float, decimals: int = 1) -> str:
    return f"{value * 100:.{decimals}f}%"


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 12,
    color: RGBColor = INK,
    bold: bool = False,
    font: str = FONT,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.02,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    p.line_spacing = 1.0
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def add_paragraphs(
    slide,
    lines: list[str],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 10.5,
    color: RGBColor = INK,
    bullet: bool = False,
    spacing: float = 4,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(0)
        p.space_after = Pt(spacing)
        p.line_spacing = 1.0
        r = p.add_run()
        r.text = ("•  " if bullet else "") + line
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return shape


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor = WHITE,
    line: RGBColor | None = LIGHT,
    radius: bool = True,
    line_width: float = 0.7,
):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width)
    return shape


def add_pill(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    *,
    fill: RGBColor = NAVY,
    color: RGBColor = WHITE,
    size: float = 7.8,
):
    add_rect(slide, x, y, w, 0.30, fill=fill, line=None, radius=True)
    add_text(
        slide,
        text,
        x + 0.04,
        y + 0.01,
        w - 0.08,
        0.28,
        size=size,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0,
    )


def add_title(slide, title: str, subtitle: str, source: str, number: int) -> None:
    add_text(slide, f"XI · {number:02d}", 0.48, 0.24, 0.62, 0.18, size=6.8, color=CYAN_DARK, bold=True)
    add_text(slide, title, 0.48, 0.49, 8.95, 0.43, size=20.5, color=NAVY, bold=True, font=FONT_MEDIUM)
    add_text(slide, subtitle, 0.49, 0.95, 8.95, 0.30, size=8.7, color=MUTED)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.49), Inches(1.27), Inches(0.74), Inches(0.035))
    line.fill.solid()
    line.fill.fore_color.rgb = CYAN
    line.line.fill.background()
    add_text(slide, f"Fuente: Informe SITFEN v12 · Cap. XI · {source}", 0.50, 5.30, 7.8, 0.13, size=5.8, color=MUTED, margin=0)
    add_text(slide, "SITFEN  |  EVALUACIÓN FINANCIERA", 8.15, 5.28, 1.35, 0.15, size=5.6, color=CYAN_DARK, bold=True, align=PP_ALIGN.RIGHT, margin=0)


def add_kpi(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    label: str,
    value: str,
    note: str = "",
    accent: RGBColor = CYAN_DARK,
    fill: RGBColor = WHITE,
    value_size: float = 19,
):
    add_rect(slide, x, y, w, h, fill=fill, line=LIGHT, radius=True)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.055), Inches(h))
    band.fill.solid()
    band.fill.fore_color.rgb = accent
    band.line.fill.background()
    add_text(slide, label.upper(), x + 0.15, y + 0.11, w - 0.25, 0.18, size=6.7, color=accent, bold=True)
    add_text(slide, value, x + 0.15, y + 0.31, w - 0.25, 0.43, size=value_size, color=NAVY, bold=True, valign=MSO_ANCHOR.MIDDLE)
    if note:
        add_text(slide, note, x + 0.15, y + h - 0.29, w - 0.25, 0.18, size=6.7, color=MUTED)


def add_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    eyebrow: str,
    title: str,
    body: str,
    accent: RGBColor = CYAN_DARK,
    fill: RGBColor = WHITE,
    title_size: float = 12.5,
    body_size: float = 8.4,
):
    add_rect(slide, x, y, w, h, fill=fill, line=LIGHT, radius=True)
    add_text(slide, eyebrow.upper(), x + 0.16, y + 0.12, w - 0.30, 0.17, size=6.5, color=accent, bold=True)
    add_text(slide, title, x + 0.16, y + 0.34, w - 0.30, 0.38, size=title_size, color=NAVY, bold=True, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, body, x + 0.16, y + 0.78, w - 0.30, h - 0.89, size=body_size, color=INK)


def add_picture_contain(slide, path: Path, x: float, y: float, w: float, h: float):
    with Image.open(path) as im:
        iw, ih = im.size
    ratio = min(w / iw, h / ih)
    pw, ph = iw * ratio, ih * ratio
    return slide.shapes.add_picture(
        str(path),
        Inches(x + (w - pw) / 2),
        Inches(y + (h - ph) / 2),
        width=Inches(pw),
        height=Inches(ph),
    )


def add_arrow(slide, x: float, y: float, w: float = 0.46, h: float = 0.24):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = CYAN_DARK
    shp.line.fill.background()
    return shp


def add_notes(slide, text: str) -> None:
    try:
        frame = slide.notes_slide.notes_text_frame
        frame.text = text
    except Exception:
        pass


def blank_layout(prs: Presentation):
    for layout in prs.slide_layouts:
        if str(layout.name).lower() == "blank":
            return layout
    return prs.slide_layouts[-1]


def new_slide(prs: Presentation):
    slide = prs.slides.add_slide(blank_layout(prs))
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = OFFWHITE
    return slide


def _matplotlib_style() -> None:
    plt.rcParams.update(
        {
            "font.family": "DejaVu Sans",
            "font.size": 10,
            "axes.titlesize": 12,
            "axes.labelcolor": "#647383",
            "xtick.color": "#647383",
            "ytick.color": "#647383",
            "axes.edgecolor": "#D6DEE5",
            "grid.color": "#E6EBEF",
            "figure.facecolor": "none",
            "axes.facecolor": "none",
            "savefig.facecolor": "none",
        }
    )


def create_charts(data: dict) -> dict[str, Path]:
    CHART_DIR.mkdir(parents=True, exist_ok=True)
    _matplotlib_style()
    base = data["base"]
    years = np.arange(1, 6)

    trajectory = CHART_DIR / "trajectory.png"
    fig, ax = plt.subplots(figsize=(8.4, 4.25), dpi=180)
    ax.plot(years, base["revenue"], marker="o", lw=3.0, color="#0B61B9", label="Ingresos")
    ax.plot(years, base["opex_total"], marker="o", lw=2.5, color="#647383", label="OPEX")
    ax.bar(years, base["ebitda"], width=0.46, color=["#C83942" if v < 0 else "#0A9196" for v in base["ebitda"]], alpha=0.88, label="EBITDA")
    ax.axhline(0, color="#AEB9C3", lw=1)
    ax.set_xticks(years, [f"Año {y}" for y in years])
    ax.set_ylabel("US$")
    ax.grid(axis="y", linewidth=0.8)
    ax.spines[["top", "right"]].set_visible(False)
    ax.legend(loc="upper left", frameon=False, ncol=3)
    for y, value in zip(years, base["revenue"]):
        ax.annotate(f"{value/1000:.1f}k", (y, value), xytext=(0, 7), textcoords="offset points", ha="center", fontsize=8, color="#0B61B9", fontweight="bold")
    fig.tight_layout(pad=0.5)
    fig.savefig(trajectory, bbox_inches="tight", transparent=True)
    plt.close(fig)

    cash = CHART_DIR / "cash_payback.png"
    fcf = np.array(base["fcf_full"], dtype=float)
    cumulative = np.cumsum(fcf)
    years0 = np.arange(0, 6)
    fig, ax = plt.subplots(figsize=(8.4, 4.1), dpi=180)
    ax.bar(years0, fcf, width=0.55, color=["#082345"] + ["#C83942" if v < 0 else "#0A9196" for v in fcf[1:]], alpha=0.92, label="FCF anual")
    ax.plot(years0, cumulative, color="#D69200", marker="o", lw=2.7, label="Caja acumulada")
    ax.axhline(0, color="#26313D", lw=1.0)
    ax.axvline(base["payback"], color="#2F7D45", ls="--", lw=1.7)
    ax.text(base["payback"] - 0.05, max(fcf) * 0.80, f"Payback {base['payback']:.2f} años", color="#2F7D45", ha="right", fontsize=9, fontweight="bold")
    ax.set_xticks(years0, ["CAPEX", "Año 1", "Año 2", "Año 3", "Año 4", "Año 5"])
    ax.set_ylabel("US$")
    ax.grid(axis="y", linewidth=0.8)
    ax.spines[["top", "right"]].set_visible(False)
    ax.legend(loc="upper left", frameon=False, ncol=2)
    fig.tight_layout(pad=0.5)
    fig.savefig(cash, bbox_inches="tight", transparent=True)
    plt.close(fig)

    tornado = CHART_DIR / "sensitivity.png"
    labels_map = {
        "Nro de gremios (adopcion)": "Adopción",
        "Precio licencia/gremio": "Precio licencia",
        "OPEX fijo (soporte/monitoreo/mantenim.)": "OPEX fijo",
        "WACC": "WACC",
        "Episodio El NiÃ±o severo (ingresos AÃ±o 2 -25%/+10%)": "El Niño severo",
        "Costo IA por consulta (Flash/Pro/voz)": "Costo IA",
    }
    rows = []
    for key, (low, high) in data["tornado"].items():
        rows.append((labels_map.get(key, key), low, high, high - low))
    rows.sort(key=lambda r: r[3])
    fig, ax = plt.subplots(figsize=(8.1, 4.25), dpi=180)
    y = np.arange(len(rows))
    for i, (label, low, high, span) in enumerate(rows):
        ax.plot([low, high], [i, i], color="#0B61B9" if label == "Adopción" else "#7E8B97", lw=10, solid_capstyle="round")
        ax.scatter([low, high], [i, i], s=18, color="#082345", zorder=3)
        ax.text(high + 900, i, f"rango {span/1000:.1f}k", va="center", fontsize=8, color="#647383")
    ax.axvline(data["base"]["npv"], color="#D69200", lw=2, ls="--", label=f"VPN base {data['base']['npv']/1000:.1f}k")
    ax.axvline(0, color="#2F7D45", lw=1.1)
    ax.set_yticks(y, [r[0] for r in rows])
    ax.set_xlabel("VPN a 5 años (US$)")
    ax.grid(axis="x", linewidth=0.8)
    ax.spines[["top", "right", "left"]].set_visible(False)
    ax.legend(loc="lower right", frameon=False)
    fig.tight_layout(pad=0.6)
    fig.savefig(tornado, bbox_inches="tight", transparent=True)
    plt.close(fig)

    mc_path = CHART_DIR / "montecarlo.png"
    mc = data["montecarlo"]
    sample = np.asarray(mc["npvs_sample"], dtype=float)
    fig, ax = plt.subplots(figsize=(7.4, 3.7), dpi=180)
    ax.hist(sample, bins=42, color="#0B61B9", alpha=0.82, edgecolor="white", linewidth=0.4)
    ax.axvline(0, color="#2F7D45", lw=2.0, label="VPN = 0")
    ax.axvline(mc["npv_p50"], color="#D69200", lw=2.0, ls="--", label=f"P50 {mc['npv_p50']/1000:.1f}k")
    ax.set_xlabel("VPN a 5 años (US$)")
    ax.set_ylabel("Frecuencia (muestra)")
    ax.grid(axis="y", linewidth=0.8)
    ax.spines[["top", "right"]].set_visible(False)
    ax.legend(loc="upper left", frameon=False)
    fig.tight_layout(pad=0.5)
    fig.savefig(mc_path, bbox_inches="tight", transparent=True)
    plt.close(fig)

    return {"trajectory": trajectory, "cash": cash, "tornado": tornado, "mc": mc_path}


def add_cover_slide(prs: Presentation, data: dict) -> None:
    slide = prs.slides.add_slide(blank_layout(prs))
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY
    add_pill(slide, "CAPÍTULO XI · CASO DE INVERSIÓN", 0.62, 0.48, 2.35, fill=CYAN_DARK)
    add_text(slide, "SITFEN", 0.62, 1.12, 4.0, 0.56, size=30, color=WHITE, bold=True, font=FONT_MEDIUM)
    add_text(slide, "Viabilidad comercial con disciplina financiera", 0.62, 1.74, 7.9, 0.64, size=24, color=WHITE, bold=True, font=FONT_MEDIUM)
    add_text(slide, "El caso no pide fe en la tecnología: pide validar adopción y liberar capital por hitos.", 0.64, 2.47, 7.55, 0.38, size=11.3, color=rgb("C7D7E7"))

    add_rect(slide, 0.63, 3.10, 3.15, 1.14, fill=AMBER, line=None, radius=True)
    add_text(slide, "GO CONDICIONAL", 0.82, 3.33, 2.78, 0.35, size=19, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "Inversión por etapas, atada a hitos comerciales", 0.82, 3.73, 2.78, 0.20, size=7.2, color=WHITE, align=PP_ALIGN.CENTER)

    add_kpi(slide, 4.10, 3.10, 1.67, 1.14, label="CAPEX", value=money(data["capex0"]), note="MVP licenciado", accent=CYAN, fill=WHITE, value_size=15)
    add_kpi(slide, 5.92, 3.10, 1.67, 1.14, label="Break-even", value="≈7 gremios", note="estructura año 5", accent=CYAN, fill=WHITE, value_size=15)
    add_kpi(slide, 7.74, 3.10, 1.67, 1.14, label="P(VPN>0) · 8a", value=pct(data["montecarlo"]["p_npv_positive_8y"], 0), note="50,000 iteraciones", accent=CYAN, fill=WHITE, value_size=18)

    add_pill(slide, "SaaS B2B2C", 0.63, 4.72, 1.18, fill=rgb("17456F"), size=7.2)
    add_pill(slide, "WACC 16%", 1.92, 4.72, 1.02, fill=rgb("17456F"), size=7.2)
    add_pill(slide, "5 años + extensión a 8", 3.05, 4.72, 1.72, fill=rgb("17456F"), size=7.2)
    add_pill(slide, "Riesgo climático NOAA + ENFEN", 4.88, 4.72, 2.35, fill=rgb("17456F"), size=7.2)
    add_text(slide, "Fuente: Informe SITFEN v12 · Capítulo XI · Tablas 16–20", 0.64, 5.25, 5.6, 0.13, size=5.8, color=rgb("8FA7BD"), margin=0)
    add_text(slide, "SITFEN  |  DEFENSA ANTE JURADO", 7.58, 5.22, 1.82, 0.15, size=5.7, color=CYAN, bold=True, align=PP_ALIGN.RIGHT, margin=0)
    add_notes(slide, "Abrir con la conclusión: GO condicional. La tecnología está construida y validada; el capital debe liberarse cuando la adopción y la disposición de pago se confirmen con gremios reales.")


def add_business_model_slide(prs: Presentation, data: dict) -> None:
    slide = new_slide(prs)
    add_title(slide, "Un contrato, muchos usuarios: SaaS B2B2C", "La unidad de facturación es el gremio; la unidad de valor es la MYPE afiliada", "Tabla 17", 2)

    add_card(slide, 0.50, 1.55, 2.24, 1.72, eyebrow="Cliente pagador", title="Gremios / asociaciones", body="Licencia anual que habilita el servicio a sus afiliados y centraliza adopción, soporte y renovación.", accent=BLUE)
    add_arrow(slide, 2.83, 2.24)
    add_card(slide, 3.40, 1.55, 2.24, 1.72, eyebrow="Producto", title="SITFEN por WhatsApp", body="Asistente agéntico con RAG, clima, voz y canalización institucional; infraestructura serverless.", accent=CYAN_DARK)
    add_arrow(slide, 5.73, 2.24)
    add_card(slide, 6.31, 1.55, 2.24, 1.72, eyebrow="Usuario final", title="20–32 MYPEs / gremio", body="Acceso cotidiano sin instalar una nueva interfaz. Uso supuesto: 40 consultas por MYPE al mes.", accent=TEAL)

    add_kpi(slide, 0.50, 3.55, 2.05, 1.18, label="Licencia anual", value="US$ 2.5k–3.5k", note="por gremio", accent=BLUE, value_size=17)
    add_kpi(slide, 2.70, 3.55, 2.05, 1.18, label="Adopción base", value="3 → 15", note="gremios, años 1–5", accent=TEAL, value_size=19)
    add_kpi(slide, 4.90, 3.55, 2.05, 1.18, label="Uso proyectado", value="28.8k → 230.4k", note="consultas/año", accent=PURPLE, value_size=15)
    add_kpi(slide, 7.10, 3.55, 2.05, 1.18, label="Costo variable", value=money(data["costo_variable_query"], 4), note="por consulta", accent=GREEN, value_size=17)

    add_rect(slide, 0.50, 4.86, 8.65, 0.26, fill=rgb("EAF7F8"), line=rgb("B7E5E9"), radius=True)
    add_text(slide, "ALCANCE  ·  El producto licenciado es el asistente y su soporte directo; Android y el portal libre no integran el CAPEX, pero el OPEX incremental del portal sí está incluido.", 0.69, 4.91, 8.27, 0.14, size=6.7, color=CYAN_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_notes(slide, "Explicar que el canal B2B2C reduce la fricción de adquisición: un contrato habilita decenas de MYPEs. Señalar con transparencia el límite del alcance financiero.")


def add_cost_slide(prs: Presentation, data: dict) -> None:
    slide = new_slide(prs)
    add_title(slide, "CAPEX contenido; economía variable casi plana", "La arquitectura serverless concentra el riesgo en adopción, no en infraestructura", "Tablas 16–17", 3)

    capex = data["capex_breakdown"]
    grouped = [
        ("Agente RAG + orquestación", capex["Desarrollo agente (LangGraph, tools/skills, router Flash/Pro, RAG)"], BLUE),
        ("Canales, WhatsApp y voz", capex["Integracion canales (WhatsApp/Twilio, voz STT/TTS, guardrails)"], CYAN_DARK),
        ("Corpus + evaluación", capex["Ingesta y curaduria corpus documental + validacion seguridad"] + capex["Evaluacion inicial (RAGAS/DeepEval, golden dataset, HITL)"], TEAL),
        ("GCP + UI + legal", capex["Infra GCP (Cloud Run, Firestore, Secret Manager, Eventarc, IAM)"] + capex["Admin UI + app de escritorio"] + capex["Legal/compliance (Ley 29733, GDPR, TyC, privacidad)"], PURPLE),
    ]
    add_text(slide, "US$ 22,000", 0.53, 1.52, 2.6, 0.45, size=25, color=NAVY, bold=True)
    add_text(slide, "CAPEX inicial del producto licenciado", 0.54, 1.98, 3.1, 0.22, size=8.4, color=MUTED)
    max_val = max(v for _, v, _ in grouped)
    top = 2.39
    for label, value, color in grouped:
        add_text(slide, label, 0.54, top, 2.22, 0.20, size=7.7, color=INK)
        add_rect(slide, 2.77, top + 0.01, 2.42, 0.16, fill=LIGHT, line=None, radius=True)
        add_rect(slide, 2.77, top + 0.01, 2.42 * value / max_val, 0.16, fill=color, line=None, radius=True)
        add_text(slide, f"US$ {value:,.0f}", 5.27, top - 0.02, 0.74, 0.20, size=7.8, color=NAVY, bold=True, align=PP_ALIGN.RIGHT)
        top += 0.55

    add_kpi(slide, 6.27, 1.55, 1.44, 1.16, label="Variable", value=money(data["costo_variable_query"], 4), note="por consulta", accent=GREEN, value_size=14.5)
    add_kpi(slide, 7.86, 1.55, 1.54, 1.16, label="OPEX fijo", value="US$ 12.45k", note="año 1", accent=BLUE, value_size=14.5)
    add_kpi(slide, 6.27, 2.88, 1.44, 1.16, label="OPEX fijo", value="US$ 23.09k", note="año 5", accent=PURPLE, value_size=14.5)
    add_kpi(slide, 7.86, 2.88, 1.54, 1.16, label="Escala de uso", value="8.0×", note="consultas A1 → A5", accent=CYAN_DARK, value_size=16)

    add_rect(slide, 6.27, 4.23, 3.13, 0.63, fill=rgb("EAF7EE"), line=rgb("BBDCC5"), radius=True)
    add_text(slide, "La demanda crece 8×, pero el OPEX total solo 2.04×.", 6.46, 4.36, 2.75, 0.21, size=9.1, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Detalle completo del CAPEX en el anexo financiero.", 0.55, 4.93, 5.4, 0.17, size=6.8, color=MUTED)
    add_notes(slide, "El punto clave no es solo el CAPEX de 22 mil dólares: es la desacoplación entre volumen y costo. La nube escala a cero y el costo unitario de IA no gobierna el caso de negocio.")


def add_growth_slide(prs: Presentation, data: dict, charts: dict[str, Path]) -> None:
    slide = new_slide(prs)
    add_title(slide, "La palanca operativa aparece desde el Año 3", "Ingresos crecen más rápido que OPEX; el margen EBITDA madura cerca de 50%", "Tabla 19 · Figura 14", 4)
    add_picture_contain(slide, charts["trajectory"], 0.45, 1.46, 6.35, 3.60)

    base = data["base"]
    margin_y5 = base["ebitda"][-1] / base["revenue"][-1]
    add_kpi(slide, 6.98, 1.53, 2.17, 1.02, label="Ingresos", value="US$ 7.5k → 52.5k", note="7.0× en 5 años", accent=BLUE, value_size=14.5)
    add_kpi(slide, 6.98, 2.70, 2.17, 1.02, label="EBITDA positivo", value="Año 3", note=f"{money(base['ebitda'][2])}", accent=TEAL, value_size=19)
    add_kpi(slide, 6.98, 3.87, 2.17, 1.02, label="Margen maduro", value=pct(margin_y5, 0), note=f"EBITDA año 5: {money(base['ebitda'][-1])}", accent=GREEN, value_size=20)
    add_text(slide, "Motor de crecimiento", 0.61, 4.88, 1.22, 0.15, size=6.2, color=CYAN_DARK, bold=True)
    add_text(slide, "3 → 15 gremios  ·  60 → 480 MYPEs  ·  28,800 → 230,400 consultas/año", 1.81, 4.86, 4.90, 0.19, size=7.6, color=NAVY, bold=True)
    add_notes(slide, "Señalar que los dos primeros años financian la maduración comercial. El EBITDA pasa de negativo a positivo en el año 3 y llega a un margen cercano a 50% en el año 5.")


def add_valuation_slide(prs: Presentation, data: dict, charts: dict[str, Path]) -> None:
    slide = new_slide(prs)
    add_title(slide, "La caja se recupera en 4.45 años; el valor exige horizonte", "El VPN a 5 años penaliza flujos tardíos; a 8 años la tesis cruza a positivo", "Tabla 18 · Figuras 15 y 19", 5)
    add_picture_contain(slide, charts["cash"], 0.44, 1.48, 6.18, 3.50)
    base = data["base"]
    add_kpi(slide, 6.86, 1.50, 1.22, 1.18, label="VPN · 5a", value=money(base["npv"]), note="destruye valor marginal", accent=RED, value_size=14)
    add_kpi(slide, 8.18, 1.50, 1.22, 1.18, label="TIR", value=pct(base["irr"]), note="< WACC 16%", accent=AMBER, value_size=18)
    add_kpi(slide, 6.86, 2.84, 1.22, 1.18, label="VPN · 8a", value=money(data["npv8_base"], plus=True), note="determinístico", accent=GREEN, value_size=14)
    add_kpi(slide, 8.18, 2.84, 1.22, 1.18, label="Payback", value=f"{base['payback']:.2f} años", note="dentro del horizonte", accent=TEAL, value_size=15)
    add_rect(slide, 6.86, 4.21, 2.54, 0.70, fill=rgb("FFF7E6"), line=rgb("E9D29B"), radius=True)
    add_text(slide, "No es un NO-GO; es una señal para financiar la escala con hitos de adopción.", 7.04, 4.34, 2.18, 0.34, size=8.4, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
    add_notes(slide, "No ocultar el VPN negativo a cinco años. Explicar que el payback ocurre cerca del borde del horizonte y que el WACC de 16% castiga los flujos tardíos. La extensión a ocho años captura la maduración B2B2C.")


def add_scenarios_slide(prs: Presentation, data: dict) -> None:
    slide = new_slide(prs)
    add_title(slide, "El rango de resultados es amplio — y comercial", "La diferencia entre extremos alcanza ≈US$ 102 mil de VPN", "Tabla 18 · Figura 16", 6)
    scenarios = data["scenarios"]
    cards = [
        ("PESIMISTA", scenarios["pesimista"], RED, "Adopción lenta + El Niño severo en el despegue"),
        ("BASE", scenarios["base"], AMBER, "Condiciones neutras + adopción según plan"),
        ("OPTIMISTA", scenarios["optimista"], GREEN, "Adopción acelerada + costo de IA descendente"),
    ]
    xs = [0.55, 3.44, 6.33]
    for x, (label, row, color, narrative) in zip(xs, cards):
        add_rect(slide, x, 1.58, 2.57, 2.76, fill=WHITE, line=color, radius=True, line_width=1.3)
        add_pill(slide, label, x + 0.18, 1.77, 1.10, fill=color, size=7.2)
        add_text(slide, money(row["npv"], plus=True), x + 0.18, 2.23, 2.21, 0.46, size=21, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, "VPN a 5 años", x + 0.18, 2.72, 2.21, 0.17, size=6.8, color=MUTED, align=PP_ALIGN.CENTER)
        irr_text = "No alcanza" if row.get("irr") is None else pct(row["irr"])
        payback_text = "> 5 años" if row.get("payback") is None else f"{row['payback']:.2f} años"
        add_text(slide, f"TIR  {irr_text}", x + 0.25, 3.06, 0.92, 0.22, size=8.5, color=color, bold=True)
        add_text(slide, f"Payback  {payback_text}", x + 1.14, 3.06, 1.17, 0.22, size=8.5, color=color, bold=True, align=PP_ALIGN.RIGHT)
        add_text(slide, narrative, x + 0.25, 3.49, 2.07, 0.48, size=7.4, color=INK, align=PP_ALIGN.CENTER)

    spread = scenarios["optimista"]["npv"] - scenarios["pesimista"]["npv"]
    add_rect(slide, 0.55, 4.60, 8.35, 0.48, fill=NAVY, line=None, radius=True)
    add_text(slide, f"Spread de VPN: {money(spread)}", 0.78, 4.71, 2.05, 0.18, size=9.1, color=CYAN, bold=True)
    add_text(slide, "La velocidad de adopción explica más valor que cualquier optimización técnica adicional.", 2.95, 4.70, 5.66, 0.19, size=8.6, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)
    add_notes(slide, "Comparar los extremos y conectar con la sensibilidad: casi 102 mil dólares separan los escenarios, y esa distancia la gobierna la adopción. Por eso el siguiente experimento debe ser comercial.")


def add_sensitivity_slide(prs: Presentation, data: dict, charts: dict[str, Path]) -> None:
    slide = new_slide(prs)
    add_title(slide, "Riesgo #1: adopción. Riesgo #6: costo de IA", "La sensibilidad ordena dónde debe trabajar la gerencia", "Figura 17 · análisis tornado", 7)
    add_picture_contain(slide, charts["tornado"], 0.44, 1.45, 6.34, 3.61)
    spans = {k: v[1] - v[0] for k, v in data["tornado"].items()}
    adoption = next(v for k, v in spans.items() if "adopcion" in k)
    price = next(v for k, v in spans.items() if "Precio" in k)
    ai = next(v for k, v in spans.items() if "Costo IA" in k)
    nino = next(v for k, v in spans.items() if "severo" in k)
    add_kpi(slide, 6.95, 1.52, 2.20, 1.04, label="Adopción", value=money(adoption), note="rango de VPN · variable dominante", accent=BLUE, value_size=17)
    add_kpi(slide, 6.95, 2.72, 2.20, 1.04, label="Precio", value=money(price), note="segunda palanca comercial", accent=TEAL, value_size=17)
    add_kpi(slide, 6.95, 3.92, 1.03, 1.04, label="El Niño", value=money(nino), note="rango", accent=AMBER, value_size=12.5)
    add_kpi(slide, 8.12, 3.92, 1.03, 1.04, label="Costo IA", value=money(ai), note="rango", accent=GREEN, value_size=12.5)
    add_text(slide, f"Adopción mueve el VPN {adoption/ai:.0f}× más que el costo de IA.", 6.96, 5.02, 2.17, 0.17, size=7.4, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(slide, "Esta es la lámina de priorización gerencial. Adopción y precio deben recibir tiempo ejecutivo; la optimización del costo de IA tiene el menor retorno marginal sobre el VPN.")


def add_montecarlo_slide(prs: Presentation, data: dict, charts: dict[str, Path]) -> None:
    slide = new_slide(prs)
    mc = data["montecarlo"]
    add_title(slide, "50,000 futuros: el horizonte cambia la probabilidad", "Monte Carlo incorpora adopción, precio, OPEX, WACC, churn y shock de El Niño", "Tabla 20 · Figura 18", 8)
    add_picture_contain(slide, charts["mc"], 0.45, 1.53, 5.55, 3.24)

    add_rect(slide, 6.20, 1.53, 1.39, 2.72, fill=rgb("FFF3F4"), line=rgb("EBC0C4"), radius=True)
    add_text(slide, "5 AÑOS", 6.41, 1.75, 0.97, 0.20, size=7.3, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, pct(mc["p_npv_positive"]), 6.36, 2.08, 1.06, 0.45, size=23, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "P(VPN > 0)", 6.42, 2.55, 0.94, 0.17, size=6.6, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, f"P10  {money(mc['npv_p10'])}\nP50  {money(mc['npv_p50'])}\nP90  {money(mc['npv_p90'], plus=True)}", 6.36, 2.94, 1.06, 0.81, size=7.0, color=INK, align=PP_ALIGN.CENTER)
    add_pill(slide, "RIESGO ALTO · CV≈0.83", 6.34, 3.82, 1.10, fill=RED, size=5.7)

    add_rect(slide, 7.78, 1.53, 1.62, 2.72, fill=rgb("EDF8F0"), line=rgb("BBDCC5"), radius=True)
    add_text(slide, "8 AÑOS", 8.09, 1.75, 0.99, 0.20, size=7.3, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, pct(mc["p_npv_positive_8y"]), 7.98, 2.08, 1.22, 0.45, size=23, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "P(VPN > 0)", 8.08, 2.55, 1.02, 0.17, size=6.6, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, f"P10  {money(mc['npv8_p10'])}\nP50  {money(mc['npv8_p50'], plus=True)}\nP90  {money(mc['npv8_p90'], plus=True)}", 7.96, 2.94, 1.26, 0.81, size=7.0, color=INK, align=PP_ALIGN.CENTER)
    add_pill(slide, "MAYORÍA CREA VALOR", 8.01, 3.82, 1.16, fill=GREEN, size=5.9)

    add_rect(slide, 0.55, 4.72, 8.85, 0.40, fill=NAVY, line=None, radius=True)
    add_text(slide, f"Shock climático calibrado: p={pct(mc['p_nino'])}/año", 0.79, 4.82, 2.42, 0.16, size=7.9, color=CYAN, bold=True)
    add_text(slide, "NOAA 19.7% + ENFEN 18.4%  ·  impacto 5–40% de ingresos  ·  severidad por bootstrap histórico", 3.17, 4.81, 5.92, 0.17, size=7.3, color=WHITE, align=PP_ALIGN.RIGHT)
    add_notes(slide, "Aclarar que 11.4% a cinco años no invalida el proyecto: muestra que la maduración comercial queda demasiado cerca del borde. A ocho años, la probabilidad supera 63% y la mediana se vuelve positiva.")


def add_evidence_slide(prs: Presentation, data: dict) -> None:
    slide = new_slide(prs)
    add_title(slide, "La tesis no depende de una sola métrica", "Cuatro lentes financieros reducen el riesgo de una conclusión optimista", "Recomendación · Anexos E–G", 9)
    cards = [
        ("01 · INVERSIÓN", "DCF + Monte Carlo", "Responde si conviene invertir bajo incertidumbre. Resultado: GO condicional y horizonte de maduración.", BLUE),
        ("02 · CLIENTE", "Unit economics", "LTV/CAC ≈ 6.8× y payback de adquisición ≈21 meses; cada gremio puede ser rentable de captar y retener.", TEAL),
        ("03 · BENEFICIO", "Infonomics", "Cada beneficio se vincula a línea base, fórmula y fuente. La reducción de rechazos se excluye por falta de evidencia.", PURPLE),
        ("04 · ARQUITECTURA", "TCO comparado", "La alternativa 100% gestionada alinea la decisión técnica con el menor costo operativo frente a build/buy.", GREEN),
    ]
    positions = [(0.55, 1.56), (5.03, 1.56), (0.55, 3.38), (5.03, 3.38)]
    for (eyebrow, title, body, color), (x, y) in zip(cards, positions):
        add_card(slide, x, y, 4.06, 1.48, eyebrow=eyebrow, title=title, body=body, accent=color, title_size=14.0, body_size=8.5)
    add_text(slide, "Rigor que convence: el modelo no monetiza beneficios sin línea base verificable.", 1.60, 5.00, 6.90, 0.20, size=8.8, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(slide, "Esta lámina anticipa una pregunta típica del jurado: ¿todo depende del VPN? No. Unit economics, beneficios defendibles y TCO corroboran el caso desde ángulos distintos.")


def add_go_slide(prs: Presentation, data: dict) -> None:
    slide = prs.slides.add_slide(blank_layout(prs))
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY
    add_pill(slide, "DECISIÓN RECOMENDADA", 0.60, 0.44, 1.72, fill=CYAN_DARK)
    add_text(slide, "GO condicional", 0.60, 1.02, 4.15, 0.55, size=29, color=WHITE, bold=True, font=FONT_MEDIUM)
    add_text(slide, "Aprobar por etapas; detener, reprecificar o acelerar según evidencia comercial.", 0.62, 1.62, 7.80, 0.32, size=11, color=rgb("C7D7E7"))

    gates = [
        ("1", "VALIDAR", "Disposición de pago y curva de adopción con PROMPERÚ/ADEX y gremios reales.", "Precio objetivo: US$ 2.5k–3.5k/año", BLUE),
        ("2", "ESCALAR", "Asegurar el umbral operativo de ≈7 gremios hacia el Año 3–4.", "Trayectoria base: 3 → 6 → 9 gremios", TEAL),
        ("3", "PROTEGER", "Cláusulas flexibles para años con El Niño y monitoreo trimestral ENFEN.", "Convertir shock de cobranza en riesgo administrado", AMBER),
    ]
    xs = [0.62, 3.62, 6.62]
    for x, (n, title, body, metric, color) in zip(xs, gates):
        add_rect(slide, x, 2.32, 2.56, 2.06, fill=rgb("102F52"), line=rgb("31516F"), radius=True)
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.18), Inches(2.52), Inches(0.38), Inches(0.38))
        circ.fill.solid(); circ.fill.fore_color.rgb = color; circ.line.fill.background()
        add_text(slide, n, x + 0.18, 2.53, 0.38, 0.36, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)
        add_text(slide, title, x + 0.69, 2.56, 1.60, 0.25, size=9.3, color=color, bold=True)
        add_text(slide, body, x + 0.19, 3.03, 2.18, 0.63, size=8.2, color=WHITE)
        add_text(slide, metric, x + 0.19, 3.86, 2.18, 0.28, size=6.6, color=rgb("AFC4D8"), bold=True, align=PP_ALIGN.CENTER)

    add_rect(slide, 0.62, 4.72, 8.56, 0.42, fill=AMBER, line=None, radius=True)
    add_text(slide, "SOLICITUD AL JURADO  ·  Aprobar la siguiente fase como experimento comercial medible, no como apuesta tecnológica abierta.", 0.84, 4.81, 8.10, 0.18, size=8.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Fuente: Informe SITFEN v12 · Capítulo XI · Recomendación", 0.62, 5.29, 5.20, 0.13, size=5.8, color=rgb("8FA7BD"), margin=0)
    add_text(slide, "SITFEN  |  GO CONDICIONAL", 7.59, 5.26, 1.60, 0.15, size=5.7, color=CYAN, bold=True, align=PP_ALIGN.RIGHT, margin=0)
    add_notes(slide, "Cerrar con una solicitud concreta: aprobar una siguiente fase acotada, con hitos de disposición de pago, adopción y umbral de gremios. El riesgo técnico ya no es el cuello de botella.")


def add_assumptions_appendix(prs: Presentation, data: dict) -> None:
    slide = new_slide(prs)
    add_title(slide, "Anexo A · Supuestos auditables del modelo", "Parámetros explícitos; cualquier jurado puede someterlos a estrés", "Tabla 17", 11)
    rows = [
        ("Horizonte / WACC / IR", "5 años (+ extensión a 8) / 16% / 29.5%", "Proyecto tech early-stage en Perú; costos en USD"),
        ("CAPEX inicial", "US$ 22,000", "Desarrollo, canales, corpus, GCP, evaluación, UI y legal"),
        ("Adopción", "3 → 15 gremios; 20 → 32 MYPEs/gremio", "Curva B2B2C por validar con gremios reales"),
        ("Precio", "US$ 2,500 → 3,500 / gremio-año", "Escala con valor entregado"),
        ("Uso", "40 consultas/MYPE/mes", "15% voz; 20% de consultas a Gemini Pro"),
        ("Costo variable", money(data["costo_variable_query"], 4) + " / consulta", "Incluye router, clasificación, visión, voz, Twilio y Firestore"),
        ("OPEX fijo", "US$ 12,450 → 23,085 / año", "Soporte, curaduría, monitoreo, mantenimiento y GCP incremental"),
        ("El Niño", "p=19.1%/año; impacto 5–40%", "Promedio NOAA/ENFEN; severidad por bootstrap histórico"),
    ]
    x0, y0 = 0.52, 1.52
    widths = [2.02, 2.73, 3.82]
    headers = ["PARÁMETRO", "VALOR", "SUSTENTO"]
    x = x0
    for header, width in zip(headers, widths):
        add_rect(slide, x, y0, width, 0.38, fill=NAVY, line=WHITE, radius=False)
        add_text(slide, header, x + 0.10, y0 + 0.08, width - 0.20, 0.18, size=7.2, color=WHITE, bold=True)
        x += width
    y = y0 + 0.38
    row_h = 0.355
    for i, row in enumerate(rows):
        fill = WHITE if i % 2 == 0 else rgb("EDF1F4")
        x = x0
        for j, (cell, width) in enumerate(zip(row, widths)):
            add_rect(slide, x, y, width, row_h, fill=fill, line=WHITE, radius=False)
            add_text(slide, cell, x + 0.09, y + 0.07, width - 0.18, row_h - 0.09, size=6.8 if j == 2 else 7.1, color=NAVY if j < 2 else INK, bold=(j == 0), valign=MSO_ANCHOR.MIDDLE)
            x += width
        y += row_h
    add_rect(slide, 0.52, 4.80, sum(widths), 0.30, fill=rgb("FFF7E6"), line=rgb("E9D29B"), radius=False)
    add_text(slide, "Límite: aplica al piloto de arándano del norte del Perú; otros cultivos/regiones requieren recalibrar adopción y precio.", 0.72, 4.87, 8.15, 0.15, size=6.8, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
    add_notes(slide, "Usar solo si el jurado pregunta por supuestos. Subrayar que todos son explícitos y reproducibles; los comerciales están deliberadamente marcados como pendientes de validación.")


def add_projection_appendix(prs: Presentation, data: dict) -> None:
    slide = new_slide(prs)
    add_title(slide, "Anexo B · Proyección base año a año", "La tabla que reconcilia adopción, ingresos, OPEX, EBITDA y caja", "Tabla 19", 12)
    base = data["base"]
    rows = [
        ("Gremios activos", [3, 6, 9, 12, 15], "int"),
        ("MYPEs afiliadas", base["mypes"], "int"),
        ("Consultas/año", base["queries_anio"], "int"),
        ("Ingresos", base["revenue"], "money"),
        ("OPEX total", base["opex_total"], "money"),
        ("EBITDA", base["ebitda"], "money"),
        ("Flujo de caja libre", base["fcf"], "money"),
    ]
    x0, y0 = 0.52, 1.55
    widths = [2.25, 1.25, 1.25, 1.25, 1.25, 1.25]
    headers = ["MÉTRICA", "AÑO 1", "AÑO 2", "AÑO 3", "AÑO 4", "AÑO 5"]
    x = x0
    for header, width in zip(headers, widths):
        add_rect(slide, x, y0, width, 0.42, fill=NAVY, line=WHITE, radius=False)
        add_text(slide, header, x + 0.05, y0 + 0.10, width - 0.10, 0.17, size=7.2, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        x += width
    y = y0 + 0.42
    row_h = 0.39
    for i, (label, values, kind) in enumerate(rows):
        fill = WHITE if i % 2 == 0 else rgb("EDF1F4")
        x = x0
        cells = [label]
        for val in values:
            if kind == "money":
                cells.append(f"{val:,.0f}")
            else:
                cells.append(f"{int(val):,}")
        for j, (cell, width) in enumerate(zip(cells, widths)):
            add_rect(slide, x, y, width, row_h, fill=fill, line=WHITE, radius=False)
            color = INK
            if label == "EBITDA" and j > 0:
                color = RED if values[j - 1] < 0 else GREEN
            add_text(slide, cell, x + 0.06, y + 0.10, width - 0.12, 0.19, size=7.2, color=color, bold=(j == 0 or label == "EBITDA"), align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT)
            x += width
        y += row_h
    add_text(slide, "Montos financieros en US$. EBITDA positivo desde el Año 3; margen EBITDA Año 5 ≈50%.", 0.72, 4.84, 8.25, 0.18, size=7.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(slide, "Usar para responder preguntas de reconciliación. La fila de gremios permite comprobar de dónde provienen ingresos, volumen y punto de equilibrio.")


def add_method_appendix(prs: Presentation, data: dict) -> None:
    slide = new_slide(prs)
    add_title(slide, "Anexo C · Riesgo modelado, no decorativo", "Código reproducible, semilla fija y calibración climática dual", "Metodología Monte Carlo · Tabla 20", 13)
    cards = [
        ("MODELO", "Flujo de caja descontado", "Horizonte 5 años y extensión analítica a 8; WACC 16%, IR 29.5%.", BLUE),
        ("INCERTIDUMBRE", "Distribuciones triangulares", "Adopción, precio, costos, WACC y churn con mínimo, moda y máximo razonados.", TEAL),
        ("CLIMA", "NOAA + ENFEN", "p=19.1% anual; severidad 5–40% por bootstrap de magnitudes históricas ONI/ICEN.", AMBER),
        ("SIMULACIÓN", "50,000 iteraciones", "La misma semilla reproduce resultados; no depende de una corrida favorable.", PURPLE),
    ]
    positions = [(0.55, 1.56), (5.03, 1.56), (0.55, 3.32), (5.03, 3.32)]
    for (eyebrow, title, body, color), (x, y) in zip(cards, positions):
        add_card(slide, x, y, 4.06, 1.40, eyebrow=eyebrow, title=title, body=body, accent=color, title_size=13.6, body_size=8.4)
    add_rect(slide, 0.55, 4.91, 8.54, 0.25, fill=NAVY, line=None, radius=True)
    add_text(slide, "Archivos reproducibles: evaluacion_financiera/finance_model.py + model_output.json", 0.76, 4.95, 8.12, 0.13, size=6.7, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_notes(slide, "Subrayar que la simulación usa fuentes climáticas independientes y semilla fija. Las distribuciones triangulares son apropiadas porque todavía no existe una historia comercial propia.")


def add_financial_module(prs: Presentation, data: dict, charts: dict[str, Path]) -> None:
    add_cover_slide(prs, data)
    add_business_model_slide(prs, data)
    add_cost_slide(prs, data)
    add_growth_slide(prs, data, charts)
    add_valuation_slide(prs, data, charts)
    add_scenarios_slide(prs, data)
    add_sensitivity_slide(prs, data, charts)
    add_montecarlo_slide(prs, data, charts)
    add_evidence_slide(prs, data)
    add_go_slide(prs, data)
    add_assumptions_appendix(prs, data)
    add_projection_appendix(prs, data)
    add_method_appendix(prs, data)


def save_module(data: dict, charts: dict[str, Path]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    while prs.slides:
        slide_id = prs.slides._sldIdLst[0]  # noqa: SLF001
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        del prs.slides._sldIdLst[0]  # noqa: SLF001
    add_financial_module(prs, data, charts)
    prs.save(OUT_MODULE)


def save_integrated(data: dict, charts: dict[str, Path]) -> None:
    if not SOURCE_DECK.exists():
        return
    prs = Presentation(SOURCE_DECK)
    add_financial_module(prs, data, charts)
    prs.core_properties.title = "SITFEN — Presentación ante jurado especializado"
    prs.core_properties.subject = "Proyecto completo con evaluación financiera y viabilidad comercial del Capítulo XI"
    prs.core_properties.keywords = "SITFEN, IA agéntica, SaaS B2B2C, evaluación financiera, viabilidad comercial"
    prs.save(OUT_INTEGRATED)


def main() -> None:
    data = json.loads(MODEL_PATH.read_text(encoding="utf-8"))
    charts = create_charts(data)
    save_module(data, charts)
    print(f"Presentación financiera independiente: {OUT_MODULE}")


if __name__ == "__main__":
    main()
