from __future__ import annotations

import argparse
import io
import os
import zipfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


FONT = "Space Grotesk"
FONT_MEDIUM = "Space Grotesk Medium"

CYAN = RGBColor(0x00, 0xD8, 0xEE)
CYAN_DARK = RGBColor(0x00, 0x91, 0xA5)
DARK = RGBColor(0x4D, 0x4B, 0x4C)
NAVY = RGBColor(0x08, 0x23, 0x45)
BLUE = RGBColor(0x0B, 0x61, 0xB9)
TEAL = RGBColor(0x0A, 0x91, 0x96)
GREEN = RGBColor(0x2F, 0x7D, 0x45)
AMBER = RGBColor(0xD6, 0x92, 0x00)
RED = RGBColor(0xC8, 0x39, 0x42)
PURPLE = RGBColor(0x69, 0x45, 0xA5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE = RGBColor(0xF7, 0xF7, 0xF7)
LIGHT = RGBColor(0xEA, 0xEA, 0xEA)
MID = RGBColor(0xB9, 0xB9, 0xB9)
MUTED = RGBColor(0x72, 0x72, 0x72)

SLIDE_W = 10.0
SLIDE_H = 5.625


def rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor.from_string(value.upper())


def delete_slide(prs: Presentation, index: int) -> None:
    slide_id = prs.slides._sldIdLst[index]  # noqa: SLF001
    rel_id = slide_id.rId
    prs.part.drop_rel(rel_id)
    del prs.slides._sldIdLst[index]  # noqa: SLF001


def move_slide(prs: Presentation, old_index: int, new_index: int) -> None:
    slide_id_list = prs.slides._sldIdLst  # noqa: SLF001
    slide_id = slide_id_list[old_index]
    del slide_id_list[old_index]
    slide_id_list.insert(new_index, slide_id)


def remove_shape(shape) -> None:
    element = shape._element  # noqa: SLF001
    element.getparent().remove(element)


def set_run(run, *, size: float, color: RGBColor, bold: bool = False, font: str = FONT) -> None:
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 16,
    color: RGBColor = DARK,
    bold: bool = False,
    font: str = FONT,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.03,
    name: str | None = None,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        shape.name = name
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    p.line_spacing = 1.0
    r = p.add_run()
    r.text = text
    set_run(r, size=size, color=color, bold=bold, font=font)
    return shape


def add_rich_text(
    slide,
    lines: list[list[tuple[str, float, RGBColor, bool]]],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    line_spacing: float = 0.95,
    margin: float = 0.03,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    for index, line in enumerate(lines):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        p.line_spacing = line_spacing
        for text, size, color, bold in line:
            r = p.add_run()
            r.text = text
            set_run(r, size=size, color=color, bold=bold)
    return shape


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor = WHITE,
    line: RGBColor | None = LIGHT,
    radius: bool = True,
    line_width: float = 0.8,
    name: str | None = None,
):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width)
    return shape


def add_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    label: str,
    title: str,
    body: str,
    accent: RGBColor = CYAN,
    value: str | None = None,
    title_size: float = 15,
    body_size: float = 10.5,
):
    card = add_rect(slide, x, y, w, h, fill=WHITE, line=rgb("D9D9D9"), radius=True)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.075), Inches(h))
    band.fill.solid()
    band.fill.fore_color.rgb = accent
    band.line.fill.background()
    compact = h < 1.20 and value is None
    label_y = y + (0.05 if compact else 0.12)
    label_h = 0.15 if compact else 0.24
    label_size = 6.6 if compact else 7.5
    add_text(slide, label.upper(), x + 0.18, label_y, w - 0.32, label_h, size=label_size, color=accent, bold=True)
    title_y = y + (0.18 if compact else 0.34)
    if value:
        add_text(slide, value, x + 0.18, title_y, w - 0.35, 0.56, size=25, color=accent, bold=True, valign=MSO_ANCHOR.MIDDLE)
        title_y += 0.58
    title_h = 0.43 if compact else 0.50
    body_y = title_y + (0.43 if compact else 0.51)
    body_h = max(0.16, h - (body_y - y) - (0.08 if compact else 0.11))
    add_text(slide, title, x + 0.18, title_y, w - 0.35, title_h, size=title_size, color=NAVY, bold=True, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, body, x + 0.18, body_y, w - 0.35, body_h, size=body_size, color=DARK)
    return card


def add_pill(slide, text: str, x: float, y: float, w: float, *, fill: RGBColor, color: RGBColor = WHITE, size: float = 9):
    add_rect(slide, x, y, w, 0.34, fill=fill, line=None, radius=True)
    add_text(slide, text, x + 0.04, y + 0.02, w - 0.08, 0.30, size=size, color=color, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def add_number_circle(slide, number: str, x: float, y: float, *, fill: RGBColor = CYAN_DARK, diameter: float = 0.36):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(diameter), Inches(diameter))
    circle.fill.solid()
    circle.fill.fore_color.rgb = fill
    circle.line.fill.background()
    add_text(slide, number, x, y, diameter, diameter, size=10.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)
    return circle


def add_title(slide, title: str, source: str, *, size: float = 23) -> None:
    add_text(slide, title, 0.64, 0.77, 8.75, 0.48, size=size, color=DARK, bold=True, font=FONT_MEDIUM, valign=MSO_ANCHOR.MIDDLE)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.66), Inches(1.27), Inches(0.74), Inches(0.035))
    line.fill.solid()
    line.fill.fore_color.rgb = CYAN
    line.line.fill.background()
    add_source(slide, source)


def add_source(slide, source: str) -> None:
    add_text(slide, f"Fuente: Informe SITFEN v12 · {source}", 0.67, 5.18, 7.95, 0.16, size=6.5, color=MUTED, valign=MSO_ANCHOR.MIDDLE, margin=0)


def add_arrow(slide, x: float, y: float, w: float = 0.32, h: float = 0.22, *, color: RGBColor = CYAN_DARK, direction: str = "right"):
    mapping = {
        "right": MSO_SHAPE.RIGHT_ARROW,
        "left": MSO_SHAPE.LEFT_ARROW,
        "down": MSO_SHAPE.DOWN_ARROW,
        "up": MSO_SHAPE.UP_ARROW,
    }
    shape = slide.shapes.add_shape(mapping[direction], Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_line(slide, x1: float, y1: float, x2: float, y2: float, *, color: RGBColor = MID, width: float = 1.4):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def add_picture_contain(slide, path: Path, x: float, y: float, w: float, h: float):
    with Image.open(path) as image:
        image_ratio = image.width / image.height
    box_ratio = w / h
    if image_ratio >= box_ratio:
        draw_w = w
        draw_h = w / image_ratio
        draw_x = x
        draw_y = y + (h - draw_h) / 2
    else:
        draw_h = h
        draw_w = h * image_ratio
        draw_y = y
        draw_x = x + (w - draw_w) / 2
    return slide.shapes.add_picture(str(path), Inches(draw_x), Inches(draw_y), width=Inches(draw_w), height=Inches(draw_h))


def extract_docx_media(docx_path: Path, media_dir: Path) -> dict[int, Path]:
    media_dir.mkdir(parents=True, exist_ok=True)
    extracted: dict[int, Path] = {}
    with zipfile.ZipFile(docx_path) as archive:
        names = set(archive.namelist())
        for index in range(2, 13):
            for ext in ("png", "jpeg", "jpg"):
                member = f"word/media/image{index}.{ext}"
                if member in names:
                    blob = archive.read(member)
                    # Some JPEGs embedded by Word contain metadata that PowerPoint
                    # rejects after python-pptx repackages them. A lossless RGB PNG
                    # normalization preserves the report figure and produces a
                    # consistently valid Office image part.
                    if ext in {"jpeg", "jpg"}:
                        output = media_dir / f"image{index}.png"
                        with Image.open(io.BytesIO(blob)) as image:
                            image.convert("RGB").save(output, format="PNG", optimize=True)
                    else:
                        output = media_dir / f"image{index}.{ext}"
                        output.write_bytes(blob)
                    extracted[index] = output
                    break
    return extracted


def replace_intro_text(prs: Presentation) -> None:
    # Slide 2 — title.
    slide = prs.slides[1]
    label = slide.shapes[0]
    label.text_frame.clear()
    p = label.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "MAESTRÍA"
    set_run(r, size=20, color=DARK, font=FONT_MEDIUM)
    title = slide.shapes[1]
    title.text_frame.clear()
    title.text_frame.word_wrap = True
    title.text_frame.margin_left = 0
    title.text_frame.margin_right = 0
    lines = [
        ("SITFEN", 47, WHITE, False),
        ("INTELIGENCIA TÉCNICA", 27, CYAN, False),
        ("Y COMERCIAL", 27, CYAN, False),
    ]
    for i, (text, size, color, bold) in enumerate(lines):
        p = title.text_frame.paragraphs[0] if i == 0 else title.text_frame.add_paragraph()
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        p.line_spacing = 0.88
        r = p.add_run()
        r.text = text
        set_run(r, size=size, color=color, bold=bold)
    add_text(slide, "MYPEs agrícolas del norte del Perú · Fenómeno de El Niño", 0.70, 4.35, 4.18, 0.44, size=11.5, color=WHITE)

    # Slide 3 — author and advisor.
    slide = prs.slides[2]
    slide.shapes[0].text_frame.clear()
    p = slide.shapes[0].text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "AUTOR"
    set_run(r, size=20, color=DARK, font=FONT_MEDIUM)
    title = slide.shapes[1]
    title.text_frame.clear()
    title.text_frame.word_wrap = True
    for i, (text, color) in enumerate([("JULIO CÉSAR", WHITE), ("SALGADO PARAGUAY", CYAN)]):
        p = title.text_frame.paragraphs[0] if i == 0 else title.text_frame.add_paragraph()
        p.space_after = Pt(0)
        p.line_spacing = 0.93
        r = p.add_run()
        r.text = text
        set_run(r, size=36, color=color)
    add_text(slide, "Asesor: Ángel Boris Alzamora Sánchez", 0.70, 4.28, 4.10, 0.42, size=12, color=WHITE)

    # Slide 4 — defense route.
    slide = prs.slides[3]
    slide.shapes[0].text_frame.clear()
    p = slide.shapes[0].text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "RUTA DE LA DEFENSA"
    set_run(r, size=15.5, color=DARK, font=FONT_MEDIUM)
    title = slide.shapes[1]
    title.text_frame.clear()
    title.text_frame.word_wrap = True
    route_lines = [
        ("PROBLEMA → MVP", 35, WHITE),
        ("ARQUITECTURA →", 29, CYAN),
        ("EVIDENCIA", 35, CYAN),
    ]
    for i, (text, size, color) in enumerate(route_lines):
        p = title.text_frame.paragraphs[0] if i == 0 else title.text_frame.add_paragraph()
        p.space_after = Pt(0)
        p.line_spacing = 0.86
        r = p.add_run()
        r.text = text
        set_run(r, size=size, color=color)
    add_text(slide, "Capítulos I–VIII", 0.70, 4.36, 4.10, 0.38, size=11.5, color=WHITE)

    # Slide 5 — transition into the case.
    slide = prs.slides[4]
    title = slide.shapes[1]
    title.text_frame.clear()
    p = title.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "SITFEN"
    set_run(r, size=30, color=DARK, bold=True)
    subtitle = slide.shapes[2]
    subtitle.text_frame.clear()
    p = subtitle.text_frame.paragraphs[0]
    p.space_after = Pt(0)
    r = p.add_run()
    r.text = "INFORMACIÓN OFICIAL "
    set_run(r, size=14.5, color=DARK)
    r = p.add_run()
    r.text = "EN EL CANAL COTIDIANO"
    set_run(r, size=14.5, color=DARK, bold=True)


def prepare_content_slides(prs: Presentation, total_content_slides: int):
    # Delete the unused seventh template slide, retain slide 6 as the first content slide.
    while len(prs.slides) > 6:
        delete_slide(prs, 6)
    base = prs.slides[5]
    for shape in list(base.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and shape.is_placeholder:
            continue
        remove_shape(shape)
    slides = [base]
    for _ in range(total_content_slides - 1):
        slides.append(prs.slides.add_slide(base.slide_layout))
    return slides


def build_content(prs: Presentation, media: dict[int, Path]) -> None:
    slides = prepare_content_slides(prs, 28)

    # Introduction is built on the final allocated content slide and moved after
    # the five opening slides once all other references have been populated.
    intro = slides[19]
    add_title(intro, "Introducción: liderazgo, asimetría y riesgo climático", "Cap. I · Introducción")
    add_rect(intro, 0.66, 1.42, 8.68, 0.61, fill=NAVY, line=None, radius=True)
    add_text(intro, "El norte del Perú combina liderazgo agroexportador con una exposición estructural al Fenómeno de El Niño.", 0.94, 1.49, 8.12, 0.45, size=12.8, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    intro_cards = [
        (
            0.66,
            "OPORTUNIDAD",
            "Motor agroexportador",
            "Arándano, palta, uva, mango y banano. En 2025 el arándano peruano llegó a 66 destinos y generó más de 100 mil empleos directos en campaña.",
            BLUE,
        ),
        (
            3.64,
            "BRECHA",
            "Información inaccesible",
            "Requisitos, protocolos fitosanitarios, tarifas e inteligencia de mercados están dispersos; la asesoría especializada queda fuera del alcance de una MYPE.",
            AMBER,
        ),
        (
            6.62,
            "RIESGO",
            "Decidir bajo el FEN",
            "El episodio 2023–2024 redujo la producción temprana de arándano. Anticipar variedades, cosecha, embarques y contratos puede cambiar una campaña.",
            RED,
        ),
    ]
    for x, label, heading, body, accent in intro_cards:
        add_card(intro, x, 2.28, 2.72, 2.00, label=label, title=heading, body=body, accent=accent, title_size=14.0, body_size=9.7)
    add_rect(intro, 1.09, 4.53, 7.82, 0.40, fill=CYAN, line=None, radius=True)
    add_text(intro, "SITFEN integra conocimiento oficial + IA agéntica por WhatsApp + monitoreo transversal del FEN.", 1.28, 4.55, 7.44, 0.36, size=10.8, color=DARK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # Chapter II additions are allocated at the end and later moved after the
    # problem/FEN sequence to preserve the report's section order.
    gaps = slides[20]
    add_title(gaps, "Brechas identificadas: tres vacíos simultáneos", "Cap. II · § 2.4")
    add_rect(gaps, 0.66, 1.42, 8.68, 0.58, fill=NAVY, line=None, radius=True)
    add_text(gaps, "Ninguna solución identificada combina corpus oficial peruano, RAG agéntico con citas, voz y monitoreo del FEN en un canal conversacional.", 0.92, 1.50, 8.16, 0.40, size=11.4, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_card(gaps, 0.66, 2.18, 2.72, 2.24, label="Brecha 1", title="Acceso", body="No existe asesoría conversacional asequible para las MYPEs en su canal natural.", accent=BLUE, value="01", title_size=15.5, body_size=10.2)
    add_card(gaps, 3.64, 2.18, 2.72, 2.24, label="Brecha 2", title="Confiabilidad", body="Las alternativas generalistas no citan fuentes oficiales verificables ni se acotan al corpus normativo vigente.", accent=PURPLE, value="02", title_size=15.5, body_size=9.7)
    add_card(gaps, 6.62, 2.18, 2.72, 2.24, label="Brecha 3", title="Integración climática", body="Ninguna herramienta conecta información comercial con el monitoreo oficial ENFEN/NOAA para decisiones de campaña.", accent=TEAL, value="03", title_size=14.2, body_size=9.6)
    add_rect(gaps, 1.18, 4.58, 7.64, 0.38, fill=CYAN, line=None, radius=True)
    add_text(gaps, "SITFEN cierra las tres brechas simultáneamente con costo marginal de centavos por consulta.", 1.38, 4.62, 7.24, 0.29, size=10.6, color=DARK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    stakeholders = slides[21]
    add_title(stakeholders, "Stakeholders: uso, adopción, conocimiento y operación", "Cap. II · § 2.5, Tabla 1")
    add_rect(stakeholders, 3.79, 2.20, 2.42, 1.10, fill=NAVY, line=None, radius=True)
    add_text(stakeholders, "SITFEN", 4.02, 2.40, 1.96, 0.37, size=21, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(stakeholders, "ecosistema B2B2C", 4.03, 2.82, 1.94, 0.23, size=8.8, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    stakeholder_cards = [
        (0.66, 1.44, 2.64, 1.18, "USAN", "MYPEs agroexportadoras", "Consultan requisitos, tarifas, mercados y clima/FEN.", GREEN),
        (0.66, 3.18, 2.64, 1.18, "USUARIO AMPLIADO", "Poblador de la zona norte", "Canaliza necesidades vinculadas al agro o al FEN.", TEAL),
        (3.55, 1.39, 2.90, 0.62, "LICENCIAN Y AMPLIFICAN", "Gremios y asociaciones · técnicos", "", BLUE),
        (6.70, 1.42, 2.64, 1.38, "ALIMENTAN Y RECIBEN", "Entidades públicas", "SENASA, SENAMHI/ENFEN, INDECI, MIDAGRI, RedCITE, PROMPERÚ, SUNAT y gobiernos.", AMBER),
        (6.70, 3.18, 2.64, 1.18, "OPERAN", "Curador del corpus", "Carga, valida y mantiene vigente la base de conocimiento.", PURPLE),
        (3.55, 3.58, 2.90, 0.78, "HABILITAN LA PLATAFORMA", "Google Cloud · Twilio · AccuWeather", "", CYAN_DARK),
    ]
    for x, y, w, h, label, heading, body, accent in stakeholder_cards:
        add_rect(stakeholders, x, y, w, h, fill=WHITE, line=accent, radius=True, line_width=1.2)
        add_text(stakeholders, label, x + 0.14, y + 0.12, w - 0.28, 0.20, size=7.4, color=accent, bold=True, align=PP_ALIGN.CENTER)
        add_text(stakeholders, heading, x + 0.14, y + 0.38, w - 0.28, 0.30, size=10.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if body:
            add_text(stakeholders, body, x + 0.16, y + 0.72, w - 0.32, h - 0.82, size=8.4, color=DARK, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_line(stakeholders, 3.30, 2.03, 3.79, 2.51, color=GREEN, width=1.4)
    add_line(stakeholders, 3.30, 3.77, 3.79, 3.02, color=TEAL, width=1.4)
    add_line(stakeholders, 5.00, 2.01, 5.00, 2.20, color=BLUE, width=1.4)
    add_line(stakeholders, 6.21, 2.51, 6.70, 2.11, color=AMBER, width=1.4)
    add_line(stakeholders, 6.21, 3.02, 6.70, 3.77, color=PURPLE, width=1.4)
    add_line(stakeholders, 5.00, 3.30, 5.00, 3.58, color=CYAN_DARK, width=1.4)

    # Chapter VI additions are allocated at the end and later inserted after
    # the prototype-to-MVP slide so that the narrative follows the report.
    mvp_value = slides[22]
    add_title(mvp_value, "Propuesta de valor del MVP", "Cap. VI · § 6.1")
    add_rect(mvp_value, 0.66, 1.46, 5.50, 2.94, fill=NAVY, line=None, radius=True)
    add_text(mvp_value, "PARA", 0.94, 1.70, 0.70, 0.22, size=8.5, color=CYAN, bold=True)
    add_text(mvp_value, "MYPEs agroexportadoras del norte y sus gremios", 1.64, 1.57, 4.18, 0.50, size=14.2, color=WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
    add_text(mvp_value, "QUE NECESITAN", 0.94, 2.13, 1.17, 0.22, size=8.5, color=CYAN, bold=True)
    add_text(mvp_value, "resolver dudas técnicas, comerciales, normativas y climáticas cuando se presentan, sin asumir costos de consultoría.", 2.12, 2.03, 3.70, 0.54, size=10.2, color=WHITE, valign=MSO_ANCHOR.MIDDLE)
    add_rect(mvp_value, 0.94, 2.76, 4.94, 1.27, fill=rgb("16385E"), line=CYAN_DARK, radius=True, line_width=1.2)
    add_text(mvp_value, "SITFEN OFRECE", 1.16, 2.94, 1.30, 0.24, size=8.8, color=CYAN, bold=True)
    add_text(mvp_value, "Un asesor virtual agéntico por WhatsApp", 2.40, 2.82, 3.16, 0.56, size=14.0, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(mvp_value, "Texto y voz · respuesta en segundos · idioma del usuario · fuentes oficiales citadas · monitoreo FEN integrado", 1.16, 3.42, 4.50, 0.39, size=9.4, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    contrasts = [
        (1.46, "CHATBOT GENERALISTA", "No conoce el corpus peruano ni garantiza fuentes oficiales vigentes.", PURPLE),
        (2.43, "PLATAFORMA DE PAGO", "Inteligencia comercial inaccesible para una MYPE y no conversacional.", AMBER),
        (3.40, "CONSULTORÍA TRADICIONAL", "Respuesta en 24–72 horas y costo de US$ 30–80 por consulta.", RED),
    ]
    for y, heading, body, accent in contrasts:
        add_rect(mvp_value, 6.40, y, 2.94, 0.82, fill=WHITE, line=accent, radius=True, line_width=1.1)
        add_text(mvp_value, heading, 6.58, y + 0.12, 2.58, 0.20, size=7.8, color=accent, bold=True, align=PP_ALIGN.CENTER)
        add_text(mvp_value, body, 6.58, y + 0.37, 2.58, 0.31, size=8.2, color=DARK, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_rect(mvp_value, 1.12, 4.63, 7.76, 0.34, fill=CYAN, line=None, radius=True)
    add_text(mvp_value, "Valor diferencial = inmediatez + asequibilidad + trazabilidad oficial + sensibilidad climática", 1.31, 4.67, 7.38, 0.25, size=10.2, color=DARK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    functional_scope = slides[23]
    add_title(functional_scope, "Alcance funcional del MVP", "Cap. VI · §§ 6.3 y 6.6")
    add_rect(functional_scope, 0.66, 1.42, 8.68, 0.48, fill=NAVY, line=None, radius=True)
    add_text(functional_scope, "El alcance se incluye por trazabilidad directa a una user story MUST o a un KPI del proyecto.", 0.94, 1.49, 8.12, 0.33, size=10.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    scope_cards = [
        (0.66, 2.08, "INTELIGENCIA DOCUMENTAL", "6 tools RAG sobre corpus oficial · requisitos, tarifas, mercados, resúmenes y documentos nuevos · citación determinista · respuesta ES/EN.", BLUE),
        (5.00, 2.08, "CANAL, MULTIMODALIDAD Y ESTADO", "WhatsApp con texto y voz bidireccional · interpretación de imágenes · memoria de 12 turnos + largo plazo · personalización y derecho al olvido.", PURPLE),
        (0.66, 3.43, "OPERACIÓN CON HUMANO EN EL CIRCUITO", "Admin UI e ingesta segura · registro de usuarios · derivación institucional con confirmación explícita · handoff y consola de operador en vivo.", GREEN),
        (5.00, 3.43, "ALERTA FEN Y CANALES COMPLEMENTARIOS", "AccuWeather + corpus ENFEN · alertas y re-enganche · portal SST, Android, avatar web y Electron como demostradores complementarios.", TEAL),
    ]
    for x, y, heading, body, accent in scope_cards:
        add_rect(functional_scope, x, y, 4.34, 1.13, fill=WHITE, line=accent, radius=True, line_width=1.2)
        add_text(functional_scope, heading, x + 0.19, y + 0.15, 3.96, 0.24, size=8.8, color=accent, bold=True, align=PP_ALIGN.CENTER)
        add_text(functional_scope, body, x + 0.21, y + 0.47, 3.92, 0.50, size=8.7, color=NAVY, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_rect(functional_scope, 1.02, 4.72, 7.96, 0.30, fill=rgb("FFF1F1"), line=rgb("E6BDBD"), radius=True)
    add_text(functional_scope, "Límite: sin WhatsApp Business plenamente productivo, facturación/multi-tenancy, fine-tuning ni integración del portal/app con el backend.", 1.20, 4.72, 7.60, 0.30, size=8.2, color=RED, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    acceptance = slides[24]
    add_title(acceptance, "Criterios de aceptación del MVP", "Cap. VI · § 6.7; Cap. VIII · § 8.9", size=22)
    acceptance_items = [
        (0.66, 1.46, 4.18, 0.66, "01", "E2E DE USUARIO", "Sandbox WhatsApp: cinco MUST, texto y voz, idioma correcto y fuentes citadas.", GREEN),
        (0.66, 2.25, 4.18, 0.66, "02", "INGESTA SEGURA", "Rechaza EICAR, contenido activo y duplicados; conserva la causa en la traza.", RED),
        (0.66, 3.04, 4.18, 0.66, "03", "GUARDRAILS Y PRIVACIDAD", "Neutraliza inyección; redacta PII y pseudonimiza al usuario antes de persistir.", PURPLE),
        (0.66, 3.83, 4.18, 0.66, "04", "DESPLIEGUE REPRODUCIBLE", "Scripts 00–07 idempotentes sobre un proyecto GCP limpio, con health checks.", BLUE),
        (5.16, 1.46, 4.18, 0.86, "05", "CALIDAD SOBRE EL PIPELINE REAL", "RAGAS 0.9375 / DeepEval 0.9458 en fidelidad y 0.9062 en relevancia.", TEAL),
        (5.16, 2.48, 4.18, 0.86, "06", "CONTROL DE COSTO", "Alerta mensual aprovisionada; gasto en reposo observado ≈ US$ 0.", AMBER),
        (5.16, 3.50, 4.18, 0.99, "07", "CI Y SEGURIDAD EN VERDE", "236 pruebas: 83 unitarias + 153 de integración, junto con pip-audit y gitleaks.", CYAN_DARK),
    ]
    for x, y, w, h, number, heading, body, accent in acceptance_items:
        add_rect(acceptance, x, y, w, h, fill=WHITE, line=accent, radius=True, line_width=1.1)
        add_number_circle(acceptance, number, x + 0.14, y + (h - 0.36) / 2, fill=accent, diameter=0.36)
        add_text(acceptance, heading, x + 0.62, y + 0.08, 1.55, 0.32, size=7.8, color=accent, bold=True, valign=MSO_ANCHOR.MIDDLE)
        add_text(acceptance, body, x + 2.08, y + 0.08, w - 2.24, h - 0.16, size=8.0, color=NAVY, bold=True, valign=MSO_ANCHOR.MIDDLE)
    add_rect(acceptance, 1.02, 4.69, 7.96, 0.34, fill=CYAN, line=None, radius=True)
    add_text(acceptance, "Definición de terminado: los 7 criterios fueron verificados de forma repetible, no solo declarados.", 1.20, 4.73, 7.60, 0.25, size=9.4, color=DARK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    figure_4 = slides[25]
    add_picture_contain(figure_4, media[5], 1.04, 0.77, 7.92, 4.46)
    add_source(figure_4, "Cap. VII · § 7.1, Fig. 4")

    figure_6 = slides[26]
    add_picture_contain(figure_6, media[7], 1.04, 0.77, 7.92, 4.46)
    add_source(figure_6, "Cap. VII · § 7.2, Fig. 6")

    figure_9 = slides[27]
    add_picture_contain(figure_9, media[10], 1.04, 0.77, 7.92, 4.46)
    add_source(figure_9, "Cap. VII · § 7.3, Fig. 9")

    # 7 — central problem.
    slide = slides[0]
    add_title(slide, "Descripción del problema central", "Cap. II · § 2.2")
    add_text(
        slide,
        "Las MYPEs agrícolas del norte no disponen de asesoría técnica y comercial oportuna, confiable y asequible para exportar bajo riesgo climático recurrente.",
        0.72,
        1.42,
        8.56,
        0.60,
        size=15.5,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_card(slide, 0.66, 2.18, 2.72, 1.78, label="01 · Dispersión", title="Conocimiento fragmentado", body="Normas, protocolos, tarifas y mercados viven en documentos técnicos extensos y heterogéneos.", accent=BLUE)
    add_card(slide, 3.63, 2.18, 2.72, 1.78, label="02 · Fricción", title="Consultoría lenta y costosa", body="La asesoría tradicional tarda 24–72 horas y cuesta US$ 30–80 por consulta puntual.", accent=AMBER)
    add_card(slide, 6.60, 2.18, 2.72, 1.78, label="03 · Canal", title="La interfaz no llega al campo", body="El usuario necesita móvil, voz, imagen, español y respuestas inmediatas con la fuente oficial citada.", accent=TEAL)
    add_rect(slide, 1.46, 4.23, 7.08, 0.55, fill=CYAN, line=None, radius=True)
    add_text(slide, "El Fenómeno de El Niño amplifica el costo de decidir tarde.", 1.60, 4.28, 6.80, 0.42, size=15, color=DARK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # 7 — FEN context.
    slide = slides[1]
    add_title(slide, "El FEN convierte la asimetría en riesgo operativo", "Caps. I–II · § 2.2, Fig. 1")
    add_picture_contain(slide, media[2], 0.46, 1.38, 5.47, 3.62)
    add_card(slide, 6.22, 1.43, 3.10, 1.00, label="Señal global · NOAA", title="19.7% anual", body="Frecuencia empírica de episodios El Niño moderados o superiores.", accent=BLUE, value=None, title_size=24, body_size=9.2)
    add_card(slide, 6.22, 2.54, 3.10, 1.00, label="Señal costera · ENFEN", title="18.4% anual", body="Frecuencia del El Niño costero medida con ICEN.", accent=TEAL, value=None, title_size=24, body_size=9.2)
    add_card(slide, 6.22, 3.65, 3.10, 1.12, label="Impacto observado", title="2023–2024", body="El episodio fuerte redujo la producción temprana de arándano por golpe de calor.", accent=RED, value=None, title_size=21, body_size=9.2)
    add_text(slide, "Dos índices complementarios: señal oceánica global y señal costera relevante para el norte peruano.", 6.26, 4.82, 3.05, 0.25, size=7.8, color=MUTED, align=PP_ALIGN.CENTER)

    # 8 — personas and value proposition.
    slide = slides[2]
    add_title(slide, "Precisión técnica en el canal cotidiano", "Caps. II, V–VI · §§ 2.7, 5.4 y 6.1")
    add_card(slide, 0.66, 1.46, 2.62, 2.65, label="Usuario final", title="Productor-exportador MYPE", body="1–20 ha · smartphone Android · notas de voz en español · 1–3 consultas por semana en campaña.\n\nNecesita lenguaje llano, inmediatez y fuente oficial para actuar.", accent=GREEN, title_size=14.2, body_size=10.2)
    add_card(slide, 3.47, 1.46, 2.62, 2.65, label="Usuario multiplicador", title="Gestor técnico del gremio", body="Atiende a decenas de afiliados y consulta a diario.\n\nNecesita respuestas consistentes, trazabilidad documental y continuidad entre sesiones.", accent=BLUE, title_size=14.2, body_size=10.2)
    add_rect(slide, 6.34, 1.46, 3.00, 2.65, fill=NAVY, line=None, radius=True)
    add_text(slide, "PROPUESTA DE VALOR", 6.58, 1.67, 2.52, 0.24, size=8, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "WhatsApp", 6.58, 1.98, 2.52, 0.43, size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "texto · voz · imagen", 6.58, 2.39, 2.52, 0.32, size=12, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Respuestas en segundos, en su idioma, construidas con fuentes oficiales citadas y monitoreo FEN integrado.", 6.62, 2.90, 2.44, 0.87, size=12, color=WHITE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_rect(slide, 1.22, 4.30, 7.55, 0.48, fill=OFFWHITE, line=rgb("D5D5D5"), radius=True)
    add_text(slide, "Modelo B2B2C", 1.40, 4.38, 1.30, 0.28, size=10, color=CYAN_DARK, bold=True, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "Gremio licencia", 2.80, 4.38, 1.42, 0.28, size=10.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_arrow(slide, 4.26, 4.41, 0.34, 0.18)
    add_text(slide, "MYPE consulta", 4.67, 4.38, 1.42, 0.28, size=10.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_arrow(slide, 6.13, 4.41, 0.34, 0.18)
    add_text(slide, "Confianza y uso", 6.54, 4.38, 1.80, 0.28, size=10.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

    # 9 — general objective.
    slide = slides[3]
    add_title(slide, "Objetivo general: construir, operar y verificar", "Cap. III · § 3.1")
    add_rect(slide, 0.66, 1.40, 8.68, 0.76, fill=NAVY, line=None, radius=True)
    add_text(slide, "Desarrollar y desplegar en GCP un sistema de información temprana de inteligencia técnica y comercial para las MYPEs agrícolas del norte del Perú ante el FEN.", 0.94, 1.50, 8.12, 0.55, size=12.3, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    objective_parts = [
        (0.66, "CONSTRUIR", "Chatbot de IA agéntica con RAG avanzado y el arándano como dominio piloto del corpus.", BLUE),
        (3.64, "OPERAR", "WhatsApp con texto, voz e imagen; canalización a entidades públicas y escalamiento humano.", PURPLE),
        (6.62, "VERIFICAR", "MVP formal: fidelidad ≥ 0.80, costo ≤ US$ 0.015 y monitoreo satelital del FEN.", GREEN),
    ]
    for x, heading, body, accent in objective_parts:
        add_rect(slide, x, 2.46, 2.72, 1.73, fill=WHITE, line=accent, radius=True, line_width=1.3)
        add_text(slide, heading, x + 0.18, 2.68, 2.36, 0.27, size=10, color=accent, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.21, 3.07, 2.30, 0.82, size=10.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_rect(slide, 1.29, 4.46, 7.42, 0.37, fill=CYAN, line=None, radius=True)
    add_text(slide, "Información oficial → decisión verificable → acción oportuna ante el FEN", 1.48, 4.51, 7.04, 0.26, size=11, color=DARK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # 10 — specific objectives.
    slide = slides[4]
    add_title(slide, "Cinco objetivos específicos encadenados", "Cap. III · § 3.2")
    specific_objectives = [
        ("OE-1", "DATOS", "Pipeline seguro y reproducible: malware/contenido activo, anti-poisoning, SHA-256 y embeddings multilingües en Firestore.", BLUE),
        ("OE-2", "INTELIGENCIA", "ReAct + Gemini 2.5: router Flash/Pro, 10 tools, visión, RAG multi-query, reranking y citación determinista.", PURPLE),
        ("OE-3", "EXPERIENCIA", "WhatsApp como canal principal —texto, voz e imagen—, memoria multiusuario, derecho al olvido e interfaces complementarias.", TEAL),
        ("OE-4", "CONFIANZA", "DevSecOps: guardrails, pseudonimización/PII, secretos, firma fail-closed, CI, despliegue serverless y observabilidad.", GREEN),
        ("OE-5", "EVIDENCIA", "RAGAS + DeepEval sobre golden dataset y evaluación de viabilidad con riesgo FEN calibrado por NOAA y ENFEN.", AMBER),
    ]
    y = 1.44
    for number, heading, body, accent in specific_objectives:
        add_rect(slide, 0.70, y, 8.62, 0.61, fill=WHITE, line=accent, radius=True, line_width=1.1)
        add_number_circle(slide, number.split("-")[1], 0.87, y + 0.125, fill=accent, diameter=0.36)
        add_text(slide, number, 1.33, y + 0.10, 0.58, 0.22, size=8.2, color=accent, bold=True, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, heading, 1.92, y + 0.10, 1.20, 0.22, size=9.0, color=accent, bold=True, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, body, 3.15, y + 0.075, 5.93, 0.40, size=8.7, color=NAVY, bold=True, valign=MSO_ANCHOR.MIDDLE)
        y += 0.68
    add_rect(slide, 1.46, 4.91, 7.08, 0.19, fill=OFFWHITE, line=rgb("D7D7D7"), radius=True)
    add_text(slide, "OE-1 habilita OE-2; OE-2 y OE-3 comparten el orquestador; OE-4 es transversal; OE-5 verifica la cadena.", 1.57, 4.92, 6.86, 0.16, size=6.8, color=MUTED, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)

    # 11 — KPI scorecard.
    slide = slides[5]
    add_title(slide, "Tres KPIs para tres riesgos del pipeline", "Cap. III · § 3.3; Cap. VI · § 6.5")
    add_card(slide, 0.66, 1.48, 2.72, 2.47, label="KPI 1 · Riesgo de alucinación", title="Fidelidad", body="Meta ≥ 0.80\nRAGAS 0.9375\nDeepEval 0.9458", accent=GREEN, value="0.94", title_size=15, body_size=11.5)
    add_card(slide, 3.64, 1.48, 2.72, 2.47, label="KPI 2 · Viabilidad de escala", title="Costo variable", body="Meta ≤ US$ 0.015\nResultado US$ 0.0137\nLínea base manual US$ 30–80", accent=BLUE, value="$0.0137", title_size=15, body_size=10.7)
    add_card(slide, 6.62, 1.48, 2.72, 2.47, label="KPI 3 · Pertinencia", title="Relevancia", body="Meta ≥ 0.80\nDeepEval 0.9062\nResultado: cumplido", accent=TEAL, value="0.9062", title_size=15, body_size=11.2)
    add_rect(slide, 0.91, 4.24, 8.18, 0.55, fill=rgb("F3F1E9"), line=rgb("D9D1B5"), radius=True)
    add_text(slide, "Lectura crítica", 1.07, 4.34, 1.20, 0.24, size=9, color=AMBER, bold=True, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "RAGAS relevancy = 0.6501; el informe atribuye la divergencia a la penalización de abstenciones honestas ante preguntas fuera de corpus.", 2.26, 4.31, 6.57, 0.30, size=8.8, color=DARK, valign=MSO_ANCHOR.MIDDLE)

    # 12 — methodology.
    slide = slides[6]
    add_title(slide, "CRISP-DM ordena el dato; Lean Startup cierra", "Cap. IV · § 4.1, Fig. 2")
    add_picture_contain(slide, media[3], 0.61, 1.51, 5.54, 3.21)
    add_card(slide, 6.34, 1.42, 2.98, 0.96, label="Construir", title="Incremento operable", body="Cada iteración entrega capacidad real del MVP.", accent=BLUE, title_size=13.0, body_size=8.4)
    add_card(slide, 6.34, 2.50, 2.98, 0.96, label="Medir", title="Pipeline de producción", body="RAGAS/DeepEval sobre respuestas generadas en vivo.", accent=TEAL, title_size=13.0, body_size=8.4)
    add_card(slide, 6.34, 3.58, 2.98, 0.96, label="Aprender", title="Cambiar la siguiente iteración", body="Los resultados gobiernan el alcance posterior.", accent=AMBER, title_size=12.0, body_size=8.4)
    add_rect(slide, 1.15, 4.56, 7.70, 0.30, fill=NAVY, line=None, radius=True)
    add_text(slide, "DevSecOps  ·  MLOps  ·  DataOps  ·  AIOps/SRE  ·  PDCA", 1.24, 4.57, 7.52, 0.27, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # 13 — prototype to MVP.
    slide = slides[7]
    add_title(slide, "Del prototipo local al MVP operativo", "Cap. V · §§ 5.2, 5.4–5.6")
    add_rect(slide, 0.70, 1.48, 3.48, 2.86, fill=WHITE, line=rgb("D7D7D7"), radius=True)
    add_pill(slide, "PROTOTIPO", 0.94, 1.69, 1.15, fill=MUTED, size=9)
    add_text(slide, "Lo que ya demostraba", 0.94, 2.13, 2.90, 0.32, size=15, color=NAVY, bold=True)
    add_text(slide, "• Ingesta local en dry-run\n• CLI sobre un subconjunto del corpus\n• Recuperación semántica KNN\n• Citas en pruebas informales", 0.97, 2.54, 2.85, 1.22, size=11.2, color=DARK)
    add_text(slide, "Faltaban canal real, memoria, guardrails, voz, nube y evaluación formal.", 0.96, 3.75, 2.87, 0.39, size=9.5, color=RED, bold=True)
    add_arrow(slide, 4.39, 2.77, 0.56, 0.36, color=CYAN_DARK)
    add_rect(slide, 5.14, 1.48, 4.16, 2.86, fill=NAVY, line=None, radius=True)
    add_pill(slide, "MVP", 5.40, 1.69, 0.76, fill=CYAN_DARK, size=9)
    add_text(slide, "Sistema para usuarios reales", 5.40, 2.13, 3.55, 0.32, size=15, color=WHITE, bold=True)
    add_text(slide, "• Agente ReAct + router Flash/Pro\n• WhatsApp con voz, imagen y citas\n• Memoria + derecho al olvido\n• Guardrails y privacidad\n• Admin UI + ingesta segura\n• GCP serverless + CI/CD + observabilidad\n• Evaluación automatizada", 5.42, 2.52, 3.52, 1.56, size=10.2, color=WHITE)
    add_rect(slide, 0.86, 4.55, 8.28, 0.34, fill=OFFWHITE, line=rgb("D3D3D3"), radius=True)
    add_text(slide, "Exclusiones deliberadas: WhatsApp Business productivo · facturación/multi-tenancy · fine-tuning · ingesta web automática no curada", 1.02, 4.59, 7.96, 0.24, size=8.6, color=MUTED, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # 14 — E2E control flow.
    slide = slides[8]
    add_title(slide, "Una consulta atraviesa 9 etapas controladas", "Caps. V y VII · §§ 5.4 y 7.2, Fig. 5")
    add_rect(slide, 0.68, 1.50, 2.35, 1.34, fill=WHITE, line=BLUE, radius=True)
    add_text(slide, "ENTRADA MULTIMODAL", 0.89, 1.68, 1.94, 0.26, size=9, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Texto · voz · imagen\npor WhatsApp", 0.90, 1.99, 1.92, 0.70, size=14.0, color=NAVY, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_rect(slide, 0.68, 3.04, 2.35, 1.50, fill=OFFWHITE, line=rgb("D3D3D3"), radius=True)
    add_text(slide, "ANTES DEL LLM", 0.89, 3.22, 1.94, 0.25, size=9, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Identidad pseudonimizada\nrate limit · guardrails\nidioma · contexto · router", 0.90, 3.52, 1.92, 0.75, size=10.7, color=DARK, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_picture_contain(slide, media[6], 3.30, 1.40, 3.18, 3.63)
    add_rect(slide, 6.83, 1.50, 2.46, 1.34, fill=WHITE, line=PURPLE, radius=True)
    add_text(slide, "CENTRO AGÉNTICO", 7.02, 1.68, 2.08, 0.26, size=9, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "ReAct decide qué tools\nRAG / clima invocar", 7.03, 1.98, 2.06, 0.70, size=13.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_rect(slide, 6.83, 3.04, 2.46, 1.50, fill=OFFWHITE, line=rgb("D3D3D3"), radius=True)
    add_text(slide, "DESPUÉS DEL LLM", 7.02, 3.22, 2.08, 0.25, size=9, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Guardrail de salida\nfooter determinista\nmemoria + traza", 7.03, 3.58, 2.06, 0.66, size=11, color=DARK, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # 15 — simplified architecture.
    architecture_summary = slides[9]
    slide = architecture_summary
    add_title(slide, "Arquitectura: un núcleo, múltiples canales", "Cap. VII · §§ 7.1–7.2, Fig. 4")
    # Channels.
    channel_specs = [("WhatsApp", GREEN), ("Web", BLUE), ("Consola", PURPLE), ("Android SST", TEAL)]
    x = 0.73
    for label, color in channel_specs:
        add_rect(slide, x, 1.45, 1.55, 0.48, fill=WHITE, line=color, radius=True)
        add_text(slide, label, x + 0.08, 1.53, 1.39, 0.28, size=10.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        x += 1.72
    add_rect(slide, 7.64, 1.45, 1.65, 0.48, fill=rgb("FFF5E2"), line=AMBER, radius=True)
    add_text(slide, "Automatización", 7.72, 1.53, 1.49, 0.28, size=10, color=AMBER, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_arrow(slide, 4.74, 2.00, 0.30, 0.22, color=CYAN_DARK, direction="down")
    # Ingress and core.
    add_rect(slide, 0.73, 2.25, 1.70, 1.54, fill=OFFWHITE, line=BLUE, radius=True)
    add_text(slide, "INGRESO", 0.94, 2.44, 1.28, 0.27, size=9, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Twilio webhook\nAPI web\nfirma fail-closed", 0.92, 2.83, 1.32, 0.66, size=10.7, color=NAVY, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_arrow(slide, 2.52, 2.89, 0.37, 0.24, color=CYAN_DARK)
    add_rect(slide, 3.00, 2.17, 4.02, 1.73, fill=NAVY, line=None, radius=True)
    add_text(slide, "AgentOrchestrator.answer()", 3.20, 2.36, 3.62, 0.32, size=14.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    core = [("Guardrails", TEAL), ("Router\nFlash/Pro", AMBER), ("ReAct\n10 tools", PURPLE), ("Memoria", GREEN)]
    cx = 3.18
    for label, color in core:
        add_rect(slide, cx, 2.86, 0.77, 0.63, fill=WHITE, line=color, radius=True)
        add_text(slide, label, cx + 0.04, 2.92, 0.69, 0.50, size=8.7, color=NAVY, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        cx += 0.90
    add_arrow(slide, 7.15, 2.89, 0.37, 0.24, color=CYAN_DARK)
    add_rect(slide, 7.64, 2.25, 1.65, 1.54, fill=rgb("FFF8EA"), line=AMBER, radius=True)
    add_text(slide, "EVENTOS", 7.84, 2.44, 1.25, 0.27, size=9, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Reporte diario\nalerta SST\nre-enganche", 7.82, 2.83, 1.29, 0.66, size=10.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    # Tools and data.
    add_arrow(slide, 4.75, 3.97, 0.30, 0.22, color=CYAN_DARK, direction="down")
    tools = [("6 tools RAG", BLUE), ("Clima + NOAA", TEAL), ("Derivación + handoff", AMBER), ("Visión", PURPLE)]
    x = 0.84
    for label, color in tools:
        add_rect(slide, x, 4.28, 1.84, 0.47, fill=WHITE, line=color, radius=True)
        add_text(slide, label, x + 0.06, 4.36, 1.72, 0.28, size=9.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        x += 1.98
    add_rect(slide, 8.76, 4.28, 0.54, 0.47, fill=NAVY, line=None, radius=True)
    add_text(slide, "IA", 8.80, 4.36, 0.46, 0.28, size=9.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "Vertex AI Gemini 2.5  ·  Firestore KNN + memoria  ·  GCS + Secret Manager  ·  LangSmith / Cloud Operations", 0.97, 4.82, 8.08, 0.23, size=8.3, color=MUTED, bold=True, align=PP_ALIGN.CENTER)

    # 16 — agent control domains.
    agent_control = slides[10]
    slide = agent_control
    add_title(slide, "Control, decisión, conocimiento, acción y memoria", "Cap. VII · § 7.2, Fig. 6")
    add_rect(slide, 3.55, 2.18, 2.90, 1.25, fill=NAVY, line=None, radius=True)
    add_text(slide, "AgentOrchestrator.answer()", 3.74, 2.39, 2.52, 0.32, size=12.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Percibe → cruza → decide → actúa", 3.76, 2.84, 2.48, 0.27, size=10.2, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    domains = [
        (0.68, 1.46, 2.47, 1.39, "1 · CONTROL", "Inyección / jailbreak\nPII · rate limit · idioma\noutput guardrail", TEAL),
        (3.55, 1.42, 2.90, 0.57, "2 · DECISIÓN", "Router Flash/Pro · urgencia · entidad", AMBER),
        (6.84, 1.46, 2.47, 1.39, "3 · CONOCIMIENTO", "6 tools RAG\nclima AccuWeather + NOAA\nvisión multimodal", BLUE),
        (0.68, 3.31, 2.47, 1.43, "4 · ESTADO / ACCIÓN", "Derivación · handoff\nreportes · alertas\nre-enganche", GREEN),
        (3.55, 3.66, 2.90, 1.08, "5 · MEMORIA", "12 turnos · resumen de largo plazo\nmemoria semántica · derecho al olvido", PURPLE),
        (6.84, 3.31, 2.47, 1.43, "6 · INTEGRACIONES", "SendGrid · Scheduler\nTwilio · Firestore\nobservabilidad", CYAN_DARK),
    ]
    for x, y, w, h, heading, body, color in domains:
        add_rect(slide, x, y, w, h, fill=WHITE, line=color, radius=True, line_width=1.2)
        if h < 0.80:
            add_text(slide, heading, x + 0.15, y + 0.08, w - 0.30, 0.20, size=8.6, color=color, bold=True, align=PP_ALIGN.CENTER)
            add_text(slide, body, x + 0.15, y + 0.31, w - 0.30, max(0.18, h - 0.37), size=8.7, color=NAVY, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        else:
            add_text(slide, heading, x + 0.15, y + 0.16, w - 0.30, 0.25, size=9, color=color, bold=True, align=PP_ALIGN.CENTER)
            add_text(slide, body, x + 0.15, y + 0.50, w - 0.30, h - 0.61, size=9.7, color=NAVY, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_line(slide, 3.15, 2.40, 3.55, 2.64, color=TEAL, width=1.4)
    add_line(slide, 5.00, 1.99, 5.00, 2.18, color=AMBER, width=1.4)
    add_line(slide, 6.45, 2.64, 6.84, 2.40, color=BLUE, width=1.4)
    add_line(slide, 3.15, 3.79, 3.55, 3.17, color=GREEN, width=1.4)
    add_line(slide, 5.00, 3.43, 5.00, 3.66, color=PURPLE, width=1.4)
    add_line(slide, 6.45, 3.17, 6.84, 3.79, color=CYAN_DARK, width=1.4)
    add_rect(slide, 1.64, 4.83, 6.72, 0.28, fill=rgb("FFF1F1"), line=rgb("EDC3C3"), radius=True)
    add_text(slide, "Derivación y handoff solo se ejecutan tras confirmar el texto real del turno actual: control determinista contra agencia excesiva.", 1.78, 4.84, 6.44, 0.27, size=7.5, color=RED, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # 17 — RAG pipeline.
    slide = slides[11]
    add_title(slide, "RAG avanzado: la confianza empieza antes", "Caps. VII–VIII · §§ 7.2, 8.4 y 8.10")
    stages = [
        ("1", "CARGA", "PDF oficial\ncurado", BLUE),
        ("2", "VALIDA", "magic bytes · malware\nactivo · PII · inyección", RED),
        ("3", "TROCEA", "380 tokens\nsolape 60 + cabecera", AMBER),
        ("4", "VECTOR", "embedding\nmultilingüe", TEAL),
        ("5", "RECUPERA", "2 reformulaciones\nKNN + deduplicación", PURPLE),
        ("6", "ORDENA", "reranking Gemini\ncontexto numerado", GREEN),
    ]
    x = 0.57
    for i, (num, heading, body, color) in enumerate(stages):
        add_rect(slide, x, 1.73, 1.32, 1.52, fill=WHITE, line=color, radius=True, line_width=1.2)
        add_number_circle(slide, num, x + 0.48, 1.57, fill=color, diameter=0.36)
        add_text(slide, heading, x + 0.10, 2.05, 1.12, 0.24, size=9.5, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.10, 2.41, 1.12, 0.61, size=8.7, color=NAVY, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if i < len(stages) - 1:
            add_arrow(slide, x + 1.35, 2.38, 0.24, 0.17, color=CYAN_DARK)
        x += 1.57
    add_rect(slide, 0.79, 3.59, 4.02, 1.04, fill=OFFWHITE, line=rgb("D4D4D4"), radius=True)
    add_text(slide, "INGESTA FAIL-CLOSED", 1.00, 3.78, 3.60, 0.24, size=9, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "El hash SHA-256 sobre texto normalizado evita duplicados; la validación ocurre antes de indexar.", 1.03, 4.09, 3.54, 0.36, size=9.5, color=DARK, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_rect(slide, 5.18, 3.59, 4.02, 1.04, fill=NAVY, line=None, radius=True)
    add_text(slide, "CITACIÓN DETERMINISTA", 5.39, 3.78, 3.60, 0.24, size=9, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "El footer se construye por software desde los ToolMessages reales: ninguna fuente puede ser inventada.", 5.42, 4.05, 3.54, 0.46, size=9.2, color=WHITE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # 18 — corpus figure; use the figure title as the slide title for maximum readable area.
    slide = slides[12]
    add_picture_contain(slide, media[12], 1.04, 0.77, 7.92, 4.46)
    add_source(slide, "Cap. VIII · § 8.4, Fig. 11")

    # 19 — security and privacy.
    slide = slides[13]
    add_title(slide, "Seguridad agéntica: evidencia antes de actuar", "Caps. VI–VIII · §§ 6.4, 7.2 y 8.9")
    security_cards = [
        (0.66, "ENTRADA E INGESTA", "Firma Twilio fail-closed\nMagic bytes + malware\nAnti-inyección / anti-poisoning\nTexto externo saneado", RED),
        (3.52, "IDENTIDAD Y PRIVACIDAD", "HMAC-SHA256 del teléfono\nRedacción de PII + Luhn\nAudio: retención máxima 1 día\nDerecho al olvido", TEAL),
        (6.38, "ACCIONES DEL AGENTE", "Catálogo cerrado de 10 entidades\nClasificador de urgencia\nConfirmación explícita por código\nFooter de fuentes por software", AMBER),
    ]
    for x, heading, body, color in security_cards:
        add_rect(slide, x, 1.48, 2.62, 2.28, fill=WHITE, line=color, radius=True, line_width=1.4)
        add_text(slide, heading, x + 0.18, 1.72, 2.26, 0.32, size=10.2, color=color, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, body, x + 0.20, 2.20, 2.22, 1.23, size=10.2, color=NAVY, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_rect(slide, 0.82, 4.04, 4.00, 0.67, fill=rgb("FFF0F0"), line=rgb("E6BDBD"), radius=True)
    add_text(slide, "FAIL-CLOSED", 1.04, 4.17, 1.04, 0.24, size=10, color=RED, bold=True, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "autenticidad · guardrails · ingesta", 2.10, 4.15, 2.47, 0.28, size=9.4, color=DARK, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_rect(slide, 5.18, 4.04, 4.00, 0.67, fill=rgb("ECF7EF"), line=rgb("BEDCC5"), radius=True)
    add_text(slide, "FAIL-OPEN", 5.40, 4.17, 0.95, 0.24, size=10, color=GREEN, bold=True, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "memoria / tool no crítica → respuesta parcial explícita", 6.35, 4.13, 2.57, 0.33, size=8.8, color=DARK, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "OWASP LLM01 · LLM03 · LLM06 · LLM07 · ASI01 · ASI06", 1.51, 4.83, 6.98, 0.20, size=8.2, color=MUTED, bold=True, align=PP_ALIGN.CENTER)

    # 20 — deployment.
    gcp_overview = slides[14]
    slide = gcp_overview
    add_title(slide, "GCP serverless: operar, auditar y escalar", "Cap. VII · §§ 7.1–7.4; Cap. VIII · § 8.8")
    # Main flow.
    flow = [
        (0.63, "CANALES", "WhatsApp\nWeb · Consola\nAndroid SST", GREEN),
        (2.38, "CLOUD RUN", "FastAPI\nservicio central\n+ jobs batch", BLUE),
        (4.38, "NÚCLEO IA", "LangGraph ReAct\nVertex AI\nGemini 2.5", PURPLE),
        (6.38, "DATOS", "Firestore KNN\nmemoria · usuarios\nGCS corpus/audio", AMBER),
        (8.15, "SALIDAS", "respuesta citada\nderivación\nalertas", TEAL),
    ]
    for i, (x, heading, body, color) in enumerate(flow):
        add_rect(slide, x, 1.56, 1.34, 1.58, fill=WHITE, line=color, radius=True, line_width=1.2)
        add_text(slide, heading, x + 0.10, 1.78, 1.14, 0.26, size=9.4, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.10, 2.17, 1.14, 0.70, size=8.7, color=NAVY, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if i < len(flow) - 1:
            add_arrow(slide, x + 1.42, 2.28, 0.25, 0.18, color=CYAN_DARK)
    # Platform strips.
    add_rect(slide, 0.66, 3.48, 4.16, 0.96, fill=NAVY, line=None, radius=True)
    add_text(slide, "OPERACIONES Y SEGURIDAD", 0.90, 3.67, 3.68, 0.24, size=9.2, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "IAM · Secret Manager · Logging · Monitoring · Trace · LangSmith", 0.91, 3.97, 3.66, 0.31, size=8.7, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, 5.18, 3.48, 4.16, 0.96, fill=OFFWHITE, line=rgb("D2D2D2"), radius=True)
    add_text(slide, "EVENTOS Y ENTREGA", 5.42, 3.67, 3.68, 0.24, size=9.2, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Scheduler · Eventarc · GitHub Actions · Cloud Build · scripts 00–07", 5.43, 3.97, 3.66, 0.31, size=8.5, color=NAVY, align=PP_ALIGN.CENTER)
    add_pill(slide, "Escala 0 → N", 2.05, 4.65, 1.48, fill=GREEN, size=9)
    add_pill(slide, "Reposo ≈ US$ 0", 4.26, 4.65, 1.70, fill=BLUE, size=9)
    add_pill(slide, "10–15 concurrentes MVP", 6.72, 4.65, 2.05, fill=PURPLE, size=8.7)

    # 21 — timeline figure.
    slide = slides[15]
    add_picture_contain(slide, media[11], 1.04, 0.77, 7.92, 4.46)
    add_source(slide, "Cap. VIII · § 8.3, Fig. 10")

    # 22 — model strategy and decisions.
    slide = slides[16]
    add_title(slide, "Modelos: complejidad solo donde crea valor", "Caps. VII–VIII · §§ 7.2, 8.5 y 8.10")
    add_rect(slide, 0.67, 1.48, 2.76, 2.03, fill=NAVY, line=None, radius=True)
    add_text(slide, "80%", 0.88, 1.73, 2.34, 0.54, size=31, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Gemini 2.5 Flash", 0.89, 2.34, 2.32, 0.30, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "inferencia por defecto · ≈ US$ 0.004/consulta", 0.91, 2.79, 2.28, 0.35, size=9.6, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, 3.62, 1.48, 2.76, 2.03, fill=WHITE, line=PURPLE, radius=True, line_width=1.4)
    add_text(slide, "20%", 3.83, 1.73, 2.34, 0.54, size=31, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Gemini 2.5 Pro", 3.84, 2.34, 2.32, 0.30, size=14, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "consultas complejas · juez independiente", 3.86, 2.79, 2.28, 0.35, size=9.6, color=DARK, align=PP_ALIGN.CENTER)
    add_rect(slide, 6.57, 1.48, 2.76, 2.03, fill=WHITE, line=TEAL, radius=True, line_width=1.4)
    add_text(slide, "0", 6.78, 1.73, 2.34, 0.54, size=31, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "LLM extra para enrutar", 6.79, 2.34, 2.32, 0.30, size=13.2, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "heurística auditable: verbos · multipregunta · ≥30 palabras", 6.81, 2.76, 2.28, 0.48, size=9.2, color=DARK, align=PP_ALIGN.CENTER)
    decisions = [
        ("RAG, no fine-tuning", "La fidelidad 0.94 se logró con corpus curado.", BLUE),
        ("ReAct prearmado", "recursion_limit=8 y caché por nivel de modelo.", PURPLE),
        ("Multimodal", "Gemini transcribe audio y comprende imágenes.", TEAL),
        ("Reranking Flash", "Evita alojar un cross-encoder adicional.", AMBER),
    ]
    x = 0.67
    for heading, body, color in decisions:
        add_rect(slide, x, 3.84, 2.08, 0.91, fill=OFFWHITE, line=color, radius=True)
        add_text(slide, heading, x + 0.13, 3.99, 1.82, 0.25, size=9.5, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.13, 4.31, 1.82, 0.29, size=8.5, color=DARK, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        x += 2.20

    # 23 — validation evidence.
    slide = slides[17]
    add_title(slide, "Evidencia: lógica, integración, IA, usuario y seguridad", "Cap. VIII · § 8.9", size=21.5)
    add_rect(slide, 0.66, 1.46, 2.10, 3.28, fill=NAVY, line=None, radius=True)
    add_text(slide, "236", 0.86, 1.80, 1.70, 0.78, size=42, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "pruebas automatizadas", 0.87, 2.66, 1.68, 0.52, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "83 unitarias\n+\n153 integración", 0.89, 3.38, 1.64, 0.74, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "offline · deterministas · CI", 0.89, 4.29, 1.64, 0.22, size=8.2, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    levels = [
        (3.05, 1.48, 5.96, 0.55, "5", "SEGURIDAD EN PARALELO", "pip-audit · gitleaks · EICAR · macros · duplicados", RED),
        (3.32, 2.15, 5.42, 0.55, "4", "USUARIO E2E", "WhatsApp texto/voz · ES + segundo idioma · producción", AMBER),
        (3.59, 2.82, 4.88, 0.55, "3", "CALIDAD DE IA", "RAGAS + DeepEval sobre pipeline real · golden 180 / muestra n=45", PURPLE),
        (3.86, 3.49, 4.34, 0.55, "2", "INTEGRACIÓN", "153 casos con dobles en memoria: Firestore, SendGrid, Scheduler", TEAL),
        (4.13, 4.16, 3.80, 0.55, "1", "UNIDAD", "83 casos: chunking, guardrails, router, idioma, clima, prompt", BLUE),
    ]
    for x, y, w, h, num, heading, body, color in levels:
        add_rect(slide, x, y, w, h, fill=WHITE, line=color, radius=True, line_width=1.2)
        add_number_circle(slide, num, x + 0.10, y + 0.10, fill=color, diameter=0.34)
        add_text(slide, heading, x + 0.54, y + 0.08, 1.69, 0.22, size=8.5, color=color, bold=True, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, body, x + 2.10, y + 0.06, w - 2.24, 0.42, size=8.1, color=DARK, valign=MSO_ANCHOR.MIDDLE)

    # 24 — conclusion and honest boundary.
    slide = slides[18]
    add_title(slide, "Hipótesis técnica cerrada; adopción abierta", "Caps. V–VIII · §§ 5.3, 5.8, 6.5 y 8.9")
    add_rect(slide, 0.66, 1.47, 4.05, 2.94, fill=NAVY, line=None, radius=True)
    add_text(slide, "EVIDENCIA TÉCNICA", 0.94, 1.73, 3.49, 0.28, size=10, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    evidence = [
        ("0.94", "fidelidad vs. meta 0.80"),
        ("$0.0137", "por consulta vs. tope $0.015"),
        ("200", "PDFs oficiales · 10 tools"),
        ("236", "pruebas automatizadas en verde"),
    ]
    y = 2.17
    for value, label in evidence:
        add_text(slide, value, 0.98, y, 1.25, 0.35, size=19, color=WHITE, bold=True, align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, label, 2.36, y, 1.98, 0.35, size=10.2, color=WHITE, valign=MSO_ANCHOR.MIDDLE)
        y += 0.53
    add_rect(slide, 5.03, 1.47, 4.29, 2.94, fill=WHITE, line=RED, radius=True, line_width=1.3)
    add_text(slide, "FRONTERA DE VALIDACIÓN", 5.31, 1.73, 3.73, 0.28, size=10, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "• Sandbox de Twilio, no WhatsApp Business productivo\n\n• Sin prueba de campo con una MYPE real\n\n• Disposición de pago, adopción e intensidad de uso aún no medidas", 5.34, 2.23, 3.66, 1.54, size=11, color=NAVY, bold=True)
    add_rect(slide, 1.17, 4.61, 7.66, 0.43, fill=CYAN, line=None, radius=True)
    add_text(slide, "El riesgo dominante ya no es técnico: es adopción y pago.", 1.36, 4.66, 7.28, 0.31, size=13, color=DARK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # Restore the report's narrative order after allocating additions at the
    # end of the working slide collection.
    move_slide(prs, list(prs.slides).index(intro), 5)
    move_slide(prs, list(prs.slides).index(gaps), 8)
    move_slide(prs, list(prs.slides).index(stakeholders), 9)
    move_slide(prs, list(prs.slides).index(mvp_value), 16)
    move_slide(prs, list(prs.slides).index(functional_scope), 17)
    move_slide(prs, list(prs.slides).index(acceptance), 18)
    move_slide(prs, list(prs.slides).index(architecture_summary), 19)
    move_slide(prs, list(prs.slides).index(figure_4), list(prs.slides).index(architecture_summary) + 1)
    move_slide(prs, list(prs.slides).index(figure_6), list(prs.slides).index(agent_control) + 1)
    move_slide(prs, list(prs.slides).index(figure_9), list(prs.slides).index(gcp_overview) + 1)


def main() -> None:
    parser = argparse.ArgumentParser(description="Genera la presentación SITFEN para jurado especializado en IA.")
    workspace = Path(__file__).resolve().parents[1]
    parser.add_argument("--template", type=Path, default=workspace.parent / "Exposición SITFEN.pptx")
    parser.add_argument("--report", type=Path, default=workspace / "informe_final" / "UTEC_Informe_Final_SITFEN_v12.docx")
    parser.add_argument("--output", type=Path, default=workspace / "entregables" / "Exposición_SITFEN_Jurado_IA.pptx")
    args = parser.parse_args()

    for source in (args.template, args.report):
        if not source.exists():
            raise FileNotFoundError(source)

    args.output.parent.mkdir(parents=True, exist_ok=True)
    media_dir = workspace / "_work_presentation" / "docx_media_generated"
    media = extract_docx_media(args.report, media_dir)
    required = {2, 3, 5, 6, 7, 10, 11, 12}
    missing = sorted(required.difference(media))
    if missing:
        raise RuntimeError(f"No se encontraron las imágenes requeridas dentro del Word: {missing}")

    prs = Presentation(args.template)
    replace_intro_text(prs)
    build_content(prs, media)

    # Preserve the 16:9 size defined by the template and remove personal metadata.
    prs.core_properties.title = "SITFEN — Exposición ante jurado especializado en IA"
    prs.core_properties.subject = "Capítulos I–VIII del Informe Final SITFEN v12"
    prs.core_properties.author = "Julio César Salgado Paraguay"
    prs.core_properties.last_modified_by = "Codex"
    prs.core_properties.comments = "Contenido basado exclusivamente en UTEC_Informe_Final_SITFEN_v12.docx, capítulos I–VIII."
    prs.save(args.output)
    print(os.fspath(args.output.resolve()))
    print(f"slides={len(prs.slides)} size={prs.slide_width / 914400:.2f}x{prs.slide_height / 914400:.2f} in")


if __name__ == "__main__":
    main()
