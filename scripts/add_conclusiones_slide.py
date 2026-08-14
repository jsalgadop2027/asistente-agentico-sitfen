"""Completa únicamente la diapositiva 45 de conclusiones del deck SITFEN."""
from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
INPUT = ROOT / "entregables" / "SITFEN_IA_Exposición_Jurado_Final_v2.pptx"
OUTPUT = ROOT / "entregables" / "SITFEN_IA_Exposición_Jurado_Final_v3.pptx"
SLIDE_INDEX = 44  # diapositiva 45

FONT = "Space Grotesk"
FONT_MEDIUM = "Space Grotesk Medium"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value.lstrip("#").upper())


NAVY = rgb("082345")
DARK = rgb("4D4B4C")
TEXT = rgb("26313D")
MUTED = rgb("6B6B6B")
WHITE = rgb("FFFFFF")
LIGHT = rgb("E4E7EA")
CYAN = rgb("00C8DE")
BLUE = rgb("1168C4")
TEAL = rgb("07969D")
AMBER = rgb("DE9400")
RED = rgb("D23B43")


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 12,
    color: RGBColor = TEXT,
    bold: bool = False,
    font: str = FONT,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.01,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    p.line_spacing = 1.0
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor = WHITE,
    line: RGBColor | None = LIGHT,
    radius: bool = True,
    line_width: float = 1.0,
):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width)
    return shape


def add_card(slide, x: float, y: float, w: float, h: float, *, fill: RGBColor, line: RGBColor):
    add_rect(slide, x + 0.035, y + 0.045, w, h, fill=rgb("C9CDD1"), line=None)
    return add_rect(slide, x, y, w, h, fill=fill, line=line, line_width=1.3)


def add_pill(slide, text: str, x: float, y: float, w: float, *, fill: RGBColor, color: RGBColor = WHITE):
    add_rect(slide, x, y, w, 0.28, fill=fill, line=None, radius=True)
    add_text(slide, text, x + 0.03, y + 0.01, w - 0.06, 0.25, size=6.7, color=color, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)


def add_result_card(
    slide,
    x: float,
    *,
    label: str,
    value: str,
    note: str,
    accent: RGBColor,
    dark: bool = False,
    value_size: float = 22,
):
    y, w, h = 2.22, 2.63, 1.37
    fill = NAVY if dark else WHITE
    add_card(slide, x, y, w, h, fill=fill, line=accent)
    title_color = CYAN if dark else accent
    value_color = WHITE if dark else accent
    note_color = WHITE if dark else MUTED
    add_text(slide, label.upper(), x + 0.15, y + 0.14, w - 0.30, 0.19, size=7.2, color=title_color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, value, x + 0.15, y + 0.38, w - 0.30, 0.40, size=value_size, color=value_color, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, note, x + 0.18, y + 0.84, w - 0.36, 0.34, size=7.0, color=note_color, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def add_summary_card(slide, x: float, *, label: str, body: str, accent: RGBColor):
    y, w, h = 3.84, 4.15, 0.73
    add_card(slide, x, y, w, h, fill=WHITE, line=accent)
    add_text(slide, label.upper(), x + 0.18, y + 0.14, 1.31, 0.18, size=6.9, color=accent, bold=True)
    add_text(slide, body, x + 1.43, y + 0.10, w - 1.62, 0.44, size=7.4, color=TEXT, valign=MSO_ANCHOR.MIDDLE)


def build_slide(slide) -> None:
    if len(slide.shapes):
        raise RuntimeError("La diapositiva 45 ya contiene elementos; se cancela para no sobrescribir contenido.")

    add_text(
        slide,
        "Cinco objetivos alcanzados; mercado por validar",
        0.67,
        0.88,
        8.60,
        0.43,
        size=25.0,
        color=DARK,
        bold=True,
        font=FONT_MEDIUM,
        valign=MSO_ANCHOR.MIDDLE,
    )
    underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.67), Inches(1.27), Inches(0.74), Inches(0.035))
    underline.fill.solid()
    underline.fill.fore_color.rgb = CYAN
    underline.line.fill.background()

    add_rect(slide, 0.67, 1.49, 8.58, 0.47, fill=NAVY, line=None, radius=True)
    add_text(slide, "5/5 OBJETIVOS ESPECÍFICOS ALCANZADOS", 0.91, 1.58, 4.75, 0.22, size=10.0, color=WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
    add_pill(slide, "OE-5 · ALCANCE MVP", 7.18, 1.58, 1.75, fill=AMBER)

    add_result_card(
        slide,
        0.67,
        label="Técnica",
        value="VIABLE",
        note="0.94 fidelidad · 200 PDFs\n10 tools · RAG con citas",
        accent=CYAN,
        dark=True,
    )
    add_result_card(
        slide,
        3.65,
        label="Operativa + regulatoria",
        value="VIABLE",
        note="GCP E2E · WhatsApp/voz/imagen\nPII · DevSecOps · canalización",
        accent=TEAL,
    )
    add_result_card(
        slide,
        6.63,
        label="Económica",
        value="PLAUSIBLE",
        note="$0.0137/consulta · LTV/CAC-proxy 6.8×\nGO condicional a adopción",
        accent=AMBER,
        value_size=19.5,
    )

    add_summary_card(
        slide,
        0.67,
        label="Aporte replicable",
        body="Corpus oficial + RAG citado + WhatsApp + clima + escalamiento al Estado.",
        accent=BLUE,
    )
    add_summary_card(
        slide,
        5.10,
        label="Frontera abierta",
        body="Piloto real: adopción, disposición de pago e intensidad de uso.",
        accent=RED,
    )

    add_rect(slide, 0.67, 4.77, 8.58, 0.35, fill=CYAN, line=None, radius=True)
    add_text(slide, "SITFEN demostró viabilidad; la siguiente prueba es comercial, no técnica.", 0.91, 4.84, 8.10, 0.18, size=9.2, color=DARK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "Fuente: Informe SITFEN v12 · Cap. XIV · Conclusión general", 0.67, 5.27, 5.55, 0.13, size=6.3, color=rgb("808080"), margin=0)

    try:
        slide.notes_slide.notes_text_frame.text = (
            "Conclusión sugerida: los cinco objetivos específicos fueron alcanzados con evidencia reproducible; "
            "el OE-5 se declara alcanzado en el alcance del MVP porque se construyó y ejecutó el instrumento de "
            "evaluación, pero la rentabilidad de mercado continúa condicionada a un piloto real. SITFEN demostró "
            "viabilidad técnica, operativa y regulatoria, y plausibilidad económica bajo condiciones verificables. "
            "El patrón es replicable a otros cultivos y amenazas climáticas, pero la siguiente prueba debe medir "
            "adopción y disposición de pago, no añadir complejidad técnica."
        )
    except Exception:
        pass


def main() -> None:
    if not INPUT.exists():
        raise FileNotFoundError(INPUT)
    shutil.copy2(INPUT, OUTPUT)
    prs = Presentation(OUTPUT)
    if len(prs.slides) != 51:
        raise RuntimeError(f"Se esperaban 51 diapositivas y se encontraron {len(prs.slides)}.")
    build_slide(prs.slides[SLIDE_INDEX])
    prs.core_properties.title = "SITFEN — Presentación final ante jurado"
    prs.core_properties.subject = "Presentación final con conclusiones del Capítulo XIV"
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    main()
