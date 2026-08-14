from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "entregables" / "Presentacion_Jurado_WhatsApp_IA_Agentica.pptx"
CHARTS = ROOT / "_work_presentation" / "whatsapp_jury_charts"

NAVY = RGBColor(8, 35, 69)
BLUE = RGBColor(0, 145, 165)
CYAN = RGBColor(0, 216, 238)
GREEN = RGBColor(37, 211, 102)
AMBER = RGBColor(230, 154, 38)
RED = RGBColor(202, 57, 66)
PURPLE = RGBColor(105, 69, 165)
INK = RGBColor(51, 57, 66)
MUTED = RGBColor(105, 113, 124)
LIGHT = RGBColor(231, 237, 242)
WHITE = RGBColor(255, 255, 255)
BG = RGBColor(247, 250, 252)
FONT = "Aptos"


def text(slide, value, x, y, w, h, size=16, color=INK, bold=False,
         align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.clear()
    box.text_frame.word_wrap = True
    box.text_frame.margin_left = box.text_frame.margin_right = Inches(.03)
    box.text_frame.margin_top = box.text_frame.margin_bottom = Inches(.02)
    box.text_frame.vertical_anchor = valign
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    p.space_before = p.space_after = Pt(0)
    r = p.add_run()
    r.text = value
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def shape(slide, kind, x, y, w, h, fill=WHITE, line=LIGHT, radius=True):
    if kind == "rect":
        kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(1)
    return s


def base(prs, title_value, section):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = BG
    shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, .16, NAVY, None, False)
    text(slide, section, .55, .38, .75, .35, 12, BLUE, True)
    text(slide, title_value, 1.15, .31, 11.4, .48, 23, NAVY, True)
    shape(slide, MSO_SHAPE.RECTANGLE, .57, .91, 1.05, .05, CYAN, None, False)
    text(slide, "SITFEN | WhatsApp + IA agéntica", .57, 7.12, 5.0, .18, 7, MUTED)
    return slide


def icon(slide, kind, x, y, color=BLUE, d=.58):
    shape(slide, MSO_SHAPE.OVAL, x, y, d, d, color, None)
    cx, cy = x + d/2, y + d/2
    if kind == "whatsapp":
        shape(slide, MSO_SHAPE.OVAL, x+.14, y+.12, d-.28, d-.27, WHITE, None)
        shape(slide, MSO_SHAPE.ISOSCELES_TRIANGLE, x+.12, y+.34, .15, .14, WHITE, None)
        text(slide, "W", x+.13, y+.14, d-.26, .25, 12, color, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    elif kind == "brain":
        for dx, dy in [(.15,.18),(.29,.12),(.37,.25),(.23,.35)]:
            shape(slide, MSO_SHAPE.OVAL, x+dx, y+dy, .09, .09, WHITE, None)
        for dx,dy,w,h in [(.21,.21,.14,.025),(.32,.18,.025,.14),(.26,.33,.13,.025)]:
            shape(slide, MSO_SHAPE.RECTANGLE, x+dx, y+dy, w,h,WHITE,None,False)
    elif kind == "mail":
        shape(slide, MSO_SHAPE.RECTANGLE, x+.12,y+.17,d-.24,d-.31,WHITE,None,False)
        text(slide,"V",x+.12,y+.14,d-.24,.22,13,color,True,PP_ALIGN.CENTER)
    elif kind == "bell":
        shape(slide, MSO_SHAPE.ARC, x+.15,y+.11,d-.30,d-.19,WHITE,None)
        shape(slide, MSO_SHAPE.RECTANGLE,x+.16,y+.35,d-.32,.05,WHITE,None,False)
        shape(slide, MSO_SHAPE.OVAL,x+.25,y+.42,.08,.08,WHITE,None)
    elif kind == "chart":
        for i, hh in enumerate([.13,.23,.34]):
            shape(slide, MSO_SHAPE.RECTANGLE,x+.14+i*.11,y+.43-hh,.07,hh,WHITE,None,False)
    elif kind == "route":
        shape(slide, MSO_SHAPE.DIAMOND,x+.16,y+.14,d-.32,d-.28,WHITE,None)
        text(slide,"?",x+.16,y+.17,d-.32,.2,12,color,True,PP_ALIGN.CENTER)
    elif kind == "face":
        shape(slide, MSO_SHAPE.OVAL,x+.12,y+.10,d-.24,d-.20,WHITE,None)
        text(slide,"•  •",x+.16,y+.17,d-.32,.14,8,color,True,PP_ALIGN.CENTER)
        text(slide,"⌣",x+.18,y+.29,d-.36,.16,13,color,True,PP_ALIGN.CENTER)
    return


def card(slide, x, y, w, h, title_value, body, icon_kind, accent=BLUE):
    shape(slide, "rect", x,y,w,h,WHITE,LIGHT)
    icon(slide,icon_kind,x+.18,y+.18,accent,.52)
    text(slide,title_value,x+.82,y+.19,w-.98,.33,14,NAVY,True)
    text(slide,body,x+.20,y+.78,w-.40,h-.90,10.5,MUTED)


def arrow(slide,x,y,w=.45,h=.22,color=BLUE,direction="right"):
    kinds={"right":MSO_SHAPE.RIGHT_ARROW,"down":MSO_SHAPE.DOWN_ARROW,"left":MSO_SHAPE.LEFT_ARROW}
    shape(slide,kinds[direction],x,y,w,h,color,None)


def node(slide,x,y,w,title_value,subtitle,kind,accent=BLUE):
    shape(slide,"rect",x,y,w,1.02,WHITE,LIGHT)
    icon(slide,kind,x+.14,y+.19,accent,.48)
    text(slide,title_value,x+.72,y+.18,w-.82,.27,12,NAVY,True)
    text(slide,subtitle,x+.72,y+.50,w-.82,.30,8.5,MUTED)


def picture(slide,name,x,y,w,h):
    p=CHARTS/name
    if p.exists():
        slide.shapes.add_picture(str(p),Inches(x),Inches(y),Inches(w),Inches(h))


def cover(prs):
    s=prs.slides.add_slide(prs.slide_layouts[6])
    bg=s.background.fill; bg.solid(); bg.fore_color.rgb=NAVY
    shape(s,MSO_SHAPE.OVAL,9.8,-1.0,4.6,4.6,RGBColor(10,92,126),None)
    shape(s,MSO_SHAPE.OVAL,10.8,4.8,3.2,3.2,BLUE,None)
    icon(s,"whatsapp",.75,.75,GREEN,.85)
    text(s,"EXPOSICIÓN ANTE EL JURADO",.75,1.85,7.4,.35,13,CYAN,True)
    text(s,"Consultas y derivaciones\ndesde WhatsApp",.75,2.28,8.7,1.35,30,WHITE,True)
    text(s,"Arquitecturas, ciclos operativos e IA agéntica para una atención trazable, empática y oportuna.",.78,3.95,7.5,.70,15,RGBColor(207,226,235))
    for i,(k,l,c) in enumerate([("brain","IA agéntica",PURPLE),("route","Derivación",AMBER),("bell","Alertas",RED)]):
        icon(s,k,.82+i*2.1,5.35,c,.55); text(s,l,1.48+i*2.1,5.48,1.25,.3,10,WHITE,True)
    text(s,"SITFEN",.78,6.75,2,.35,17,WHITE,True)


def consultations_arch(prs):
    s=base(prs,"Arquitectura de las consultas","01")
    nodes=[(.55,"Ciudadano","Pregunta natural","whatsapp",GREEN),(3.05,"Twilio / API","Recepción segura","route",BLUE),(5.55,"Orquestador IA","Intención + contexto","brain",PURPLE),(8.05,"RAG documental","Búsqueda sustentada","chart",BLUE),(10.55,"Respuesta","Mensaje + fuentes","whatsapp",GREEN)]
    for i,(x,t,sub,k,c) in enumerate(nodes):
        node(s,x,2.1,2.15,t,sub,k,c)
        if i<4: arrow(s,x+2.2,2.48,.70,.25)
    card(s,.75,4.25,3.65,1.65,"Comprensión","Detecta intención, idioma y contexto conversacional.","brain",PURPLE)
    card(s,4.85,4.25,3.65,1.65,"Conocimiento","Recupera evidencia del corpus y controla alucinaciones.","chart",BLUE)
    card(s,8.95,4.25,3.65,1.65,"Trazabilidad","Registra consulta, respuesta, fuentes y métricas.","route",AMBER)


def consultations_cycle(prs):
    s=base(prs,"Ciclo de consultas: preguntas y respuestas","02")
    steps=[("1","Pregunta","whatsapp",GREEN),("2","Comprender","brain",PURPLE),("3","Buscar","chart",BLUE),("4","Construir","brain",PURPLE),("5","Responder","whatsapp",GREEN),("6","Aprender","route",AMBER)]
    for i,(n,t,k,c) in enumerate(steps):
        x=.55+i*2.1
        icon(s,k,x+.55,1.75,c,.64); text(s,n,x,2.53,.32,.32,12,WHITE,True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
        shape(s,MSO_SHAPE.OVAL,x,2.50,.34,.34,c,None)
        text(s,t,x+.40,2.50,1.48,.34,12,NAVY,True)
        if i<5: arrow(s,x+1.66,2.05,.35,.18,c)
    picture(s,"consultas_por_dia.png",.65,3.35,5.8,2.85)
    picture(s,"temas_consultas.png",6.85,3.35,5.8,2.85)


def agent_arch(prs):
    s=base(prs,"Arquitectura de IA agéntica","03")
    node(s,.65,2.25,2.1,"Entrada WhatsApp","Mensaje + contexto","whatsapp",GREEN)
    node(s,10.55,2.25,2.1,"Acción / respuesta","Consulta o derivación","route",AMBER)
    shape(s,MSO_SHAPE.OVAL,5.4,1.65,2.55,2.55,NAVY,None)
    icon(s,"brain",6.25,2.05,PURPLE,.85)
    text(s,"ORQUESTADOR",5.65,3.05,2.05,.3,13,WHITE,True,PP_ALIGN.CENTER)
    agents=[(3.15,1.35,"Agente de intención","brain",PURPLE),(8.1,1.35,"Agente RAG","chart",BLUE),(3.15,4.15,"Agente emocional","face",RED),(8.1,4.15,"Agente derivador","route",AMBER)]
    for x,y,t,k,c in agents: node(s,x,y,2.25,t,"Herramienta especializada",k,c)
    arrow(s,2.8,2.62,.55,.22); arrow(s,9.95,2.62,.55,.22)
    text(s,"Planifica • selecciona herramientas • valida • registra",3.8,6.35,5.75,.42,13,NAVY,True,PP_ALIGN.CENTER)


def referrals_arch(prs):
    s=base(prs,"Arquitectura de las derivaciones (IA agéntica)","04")
    stages=[(.55,"Señal detectada","Intención, riesgo o solicitud","face",RED),(3.05,"Evaluación IA","Entidad, urgencia y emoción","brain",PURPLE),(5.55,"Reglas de decisión","Consentimiento y prioridad","route",AMBER),(8.05,"Multienvío","Correo a responsables","mail",BLUE),(10.55,"Seguimiento","Estado y recordatorios","bell",GREEN)]
    for i,(x,t,sub,k,c) in enumerate(stages):
        node(s,x,2.05,2.15,t,sub,k,c)
        if i<4: arrow(s,x+2.18,2.43,.72,.24)
    shape(s,"rect",1.1,4.25,11.1,1.55,NAVY,None)
    for i,(v,l,c) in enumerate([("ENTIDAD","Quién atiende",BLUE),("PRIORIDAD","Cuándo actuar",RED),("CONTEXTO","Qué ocurrió",PURPLE),("TRAZA","Cómo evoluciona",GREEN)]):
        x=1.45+i*2.75
        text(s,v,x,4.62,2.2,.3,12,c,True,PP_ALIGN.CENTER)
        text(s,l,x,5.03,2.2,.25,9,WHITE,False,PP_ALIGN.CENTER)


def referrals_cycle(prs):
    s=base(prs,"Ciclo de derivaciones","05")
    steps=[("Detectar","face",RED),("Clasificar","brain",PURPLE),("Confirmar","whatsapp",GREEN),("Derivar","mail",BLUE),("Atender","route",AMBER),("Cerrar","bell",GREEN)]
    for i,(t,k,c) in enumerate(steps):
        x=.75+i*2.05
        icon(s,k,x,1.60,c,.68); text(s,str(i+1),x+.18,2.42,.32,.32,10,WHITE,True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
        shape(s,MSO_SHAPE.OVAL,x+.17,2.41,.34,.34,c,None)
        text(s,t,x-.25,2.87,1.2,.28,11,NAVY,True,PP_ALIGN.CENTER)
        if i<5: arrow(s,x+.78,1.84,.85,.20,c)
    picture(s,"derivaciones_entidades.png",.65,3.45,5.8,2.85)
    picture(s,"derivaciones_urgencia.png",6.85,3.45,5.8,2.85)


def emotions(prs):
    s=base(prs,"Análisis de las emociones","06")
    picture(s,"emociones.png",.60,1.45,6.4,4.85)
    card(s,7.35,1.48,5.25,1.25,"1. Detectar","Analiza tono, palabras clave e intensidad del mensaje.","face",PURPLE)
    card(s,7.35,2.95,5.25,1.25,"2. Adaptar","Ajusta claridad, empatía y longitud de la respuesta.","whatsapp",GREEN)
    card(s,7.35,4.42,5.25,1.25,"3. Escalar","Prioriza señales de angustia, riesgo o urgencia.","bell",RED)
    text(s,"La emoción orienta la atención; no reemplaza el criterio profesional.",7.42,6.07,5.05,.48,11,NAVY,True,PP_ALIGN.CENTER)


def email(prs):
    s=base(prs,"Envío de correos: multienvío inteligente","07")
    node(s,.65,2.35,2.2,"Caso derivado","Resumen estructurado","route",AMBER)
    arrow(s,2.95,2.73,.60,.22)
    node(s,3.65,2.35,2.25,"Compositor IA","Asunto + contexto + prioridad","brain",PURPLE)
    arrow(s,6.00,2.73,.60,.22)
    icon(s,"mail",6.75,2.30,BLUE,.92)
    text(s,"MULTIENVÍO",6.42,3.32,1.6,.28,11,NAVY,True,PP_ALIGN.CENTER)
    for i,(t,c) in enumerate([("Responsable principal",BLUE),("Equipo de apoyo",GREEN),("Supervisión / copia",PURPLE)]):
        y=1.45+i*1.52; arrow(s,7.80,y+.48,.55,.20,c); node(s,8.45,y,3.9,t,"Correo trazable con ID de caso","mail",c)
    card(s,.85,4.65,5.0,1.35,"Contenido consistente","Incluye resumen, evidencia, urgencia, datos de contacto y próximos pasos.","mail",BLUE)
    card(s,6.15,5.55,6.0,.80,"Control de entrega","Registra destinatarios, fecha, estado y reintentos.","bell",GREEN)


def statistics(prs):
    s=base(prs,"Estadísticas de las atenciones","08")
    metrics=[("Consultas","Demanda y tendencia","chart",BLUE),("Derivaciones","Destino y prioridad","route",AMBER),("Emociones","Clima de atención","face",PURPLE),("Alertas","Riesgo y oportunidad","bell",RED)]
    for i,(t,b,k,c) in enumerate(metrics): card(s,.55+i*3.15,1.35,2.85,1.2,t,b,k,c)
    picture(s,"consultas_por_dia.png",.55,2.95,4.0,3.05)
    picture(s,"derivaciones_entidades.png",4.67,2.95,4.0,3.05)
    picture(s,"alertas_urgencia.png",8.79,2.95,4.0,3.05)
    text(s,"Tablero operativo para medir volumen, oportunidad, carga y resultados.",2.25,6.40,8.8,.38,13,NAVY,True,PP_ALIGN.CENTER)


def notifications(prs):
    s=base(prs,"Notificaciones automáticas","09")
    shape(s,MSO_SHAPE.OVAL,5.53,2.05,2.25,2.25,NAVY,None); icon(s,"bell",6.20,2.47,RED,.90)
    text(s,"MOTOR DE ALERTAS",5.70,3.55,1.92,.28,11,WHITE,True,PP_ALIGN.CENTER)
    items=[(.55,1.25,"Caso urgente","Alerta inmediata","bell",RED),(9.75,1.25,"Sin respuesta","Recordatorio por SLA","mail",AMBER),(.55,4.65,"Cambio de estado","Aviso al ciudadano","whatsapp",GREEN),(9.75,4.65,"Cierre de caso","Confirmación y encuesta","chart",BLUE)]
    for x,y,t,b,k,c in items:
        node(s,x,y,3.0,t,b,k,c)
        arrow(s,3.75 if x<5 else 9.05,y+.40,.70,.22,c,"right" if x<5 else "left")
    picture(s,"alertas_urgencia.png",4.15,4.75,5.05,1.65)
    text(s,"Eventos + reglas + SLA = atención proactiva y trazable",3.25,6.55,6.85,.34,12,NAVY,True,PP_ALIGN.CENTER)


def main():
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    cover(prs); consultations_arch(prs); consultations_cycle(prs); agent_arch(prs)
    referrals_arch(prs); referrals_cycle(prs); emotions(prs); email(prs)
    statistics(prs); notifications(prs)
    OUT.parent.mkdir(parents=True,exist_ok=True)
    prs.save(OUT)
    print(f"Deck generado: {OUT}")


if __name__ == "__main__":
    main()
